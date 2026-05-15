"""
Extend the existing VP review deck with new methodology slides and the
monthly category indexes table.

Reads:  inventory_forecaster_VP_review.pptx (existing deck)
        scripts/derived_category_profiles.json (latest profiles)
Writes: inventory_forecaster_VP_review.pptx (in place, with new slides appended)
"""
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

DECK_PATH    = Path("inventory_forecaster_VP_review.pptx")
PROFILE_PATH = Path("scripts/derived_category_profiles.json")

# ── colors / fonts ──
NAVY   = RGBColor(0x0A, 0x29, 0x52)
ORANGE = RGBColor(0xD9, 0x6B, 0x1F)
GREY   = RGBColor(0x55, 0x55, 0x55)
LIGHT  = RGBColor(0xF4, 0xF1, 0xEC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x1E, 0x88, 0x4F)


def add_title(slide, text, top=Inches(0.4)):
    box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.33), Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = NAVY


def add_subtitle(slide, text, top=Inches(1.05)):
    box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.33), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(14)
    r.font.italic = True
    r.font.color.rgb = GREY


def add_bullets(slide, items, left, top, width, height, font_size=14):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if isinstance(item, tuple) and len(item) == 2 and isinstance(item[0], int):
            level, text = item
            p.level = level
        else:
            text = item
        r = p.add_run()
        r.text = "• " + text if not text.startswith(("•", "—", "·")) else text
        r.font.size = Pt(font_size)
        r.font.color.rgb = NAVY
        p.space_after = Pt(4)


# ─── Build slides ───
prs = Presentation(str(DECK_PATH))
blank_layout = prs.slide_layouts[6]   # blank


# ── Slide A: New methodology stack overview ──
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Empirical Sales Index Methodology — 2026 Q2 Update")
add_subtitle(slide, "Six-stage pipeline applied to 2024–2026 invoice ship history (~738K rows, 31 cats / 85 sub-cats)")

bullets = [
    ("1. Source: direct /v1/records/query against ProductTrack.Invoices (table bpaxk2v8t), filter Shpd Date >= 2024-01-01, full customer base"),
    ("2. SKU consistency filter — only MStyles shipping in >=10 distinct months over a >=12mo lifespan with >=50% activity rate (excludes 1-off promo blasts)"),
    (1, "1,592 of 5,831 MStyles qualify (27%) — capturing 56% of unit volume"),
    ("3. Tariff-OOS exclusion — drop 2025 May-Sep entirely so artificially-suppressed shipments don't bias the seasonal shape"),
    ("4. Year weighting: 2024 = 2.0× (clean pre-tariff baseline), 2025 = 1.0×, 2026 = 1.0×"),
    ("5. Strategic-customer weighting: AMAZON / WAL MART / PETSMART = 2.0× (compounds with year weight)"),
    ("6. Holiday lead-time uplift: Sep ×1.10, Oct ×1.20, Nov ×1.15 — reflects 4–6 week shipping lead before Nov/Dec consumer demand peaks"),
    ("7. Planner overrides: hand-curated profiles for categories where data shape ≠ business reality (e.g. Disposable Tabletop summer + Thanksgiving + Holiday)"),
    ("8. Quality gates: total ≥ 100K units · ≥ 3 active months · peak/trough ≥ 1.30 · ≥ 2 contributing years · ≥ 3 consistent SKUs"),
]
add_bullets(slide, bullets, Inches(0.5), Inches(1.55), Inches(12.3), Inches(5.5), font_size=13)

# Footer note
note = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4))
note.text_frame.text = "All adjustments stack multiplicatively, then renormalize so each profile mean = 1.0; values clamped to [0.10, 4.00]"
note.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
note.text_frame.paragraphs[0].runs[0].font.italic = True
note.text_frame.paragraphs[0].runs[0].font.color.rgb = GREY


# ── Slide B: Forecasting consumption rules ──
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Forecasting Consumption Rules")
add_subtitle(slide, "How _get_category_profile() applies the saved indexes")

bullets = [
    ("RULE 1 — SKU gate: only apply seasonality if matched category has > 10 consistent SKUs"),
    (1, "Below the gate, fall through to next-priority match (subcat → cat → keyword fallback → none)"),
    ("RULE 2 — Floor at 1.00: every month value capped at max(value, 1.00)"),
    (1, "Seasonal indexes ONLY increase forecast demand — they never suppress it"),
    (1, "Protects against noisy 'down' signals biasing forecasts negatively"),
    ("RULE 3 — Match priority (first match wins):"),
    (1, "1) Explicit Season tag from Quickbase Styles.[Season]  (planner-curated)"),
    (1, "2) Empirical (Category, Subcategory) — 105 subcats with multi-year data"),
    (1, "3) Empirical Category alone — 31 categories with multi-year data"),
    (1, "4) Hand-curated keyword fallback (CATEGORY_PROFILES) — for items missing tags"),
    ("Constants in scripts/inventory_forecaster.py:"),
    (1, "SEASONAL_MIN_SKU_COUNT = 10    # require strictly > this many SKUs"),
    (1, "SEASONAL_FLOOR         = 1.00  # never multiply demand below 1.0× baseline"),
]
add_bullets(slide, bullets, Inches(0.5), Inches(1.55), Inches(12.3), Inches(5.5), font_size=13)


# ── Slide C: Monthly category indexes table ──
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Monthly Sales Indexes by Category")
add_subtitle(slide, "Mean = 1.00 · clamped [0.10, 4.00] · Sep-Nov lifted for holiday lead-time · sorted by total units")

prof = json.loads(PROFILE_PATH.read_text())
cats = prof["by_category"]
sorted_cats = sorted(cats.items(), key=lambda kv: -kv[1]["stats"]["total_units"])

months = ["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]

# Table layout
n_rows = 1 + len(sorted_cats)            # header + data rows
n_cols = 1 + 12 + 2                      # category + 12 months + #SKU + Peak
table_left = Inches(0.25)
table_top  = Inches(1.4)
table_w    = Inches(12.83)
table_h    = Inches(0.25 * n_rows + 0.25)

tbl_shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_w, table_h)
tbl = tbl_shape.table

# Column widths
col_widths_in = [2.0] + [0.62] * 12 + [0.55, 0.65]
for ci, w in enumerate(col_widths_in):
    tbl.columns[ci].width = Inches(w)
# Row heights
for ri in range(n_rows):
    tbl.rows[ri].height = Inches(0.22)

# Header row
headers = ["Category"] + months + ["#SKU", "Peak"]
for ci, h in enumerate(headers):
    cell = tbl.cell(0, ci)
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    tf = cell.text_frame
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    tf.text = h
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = WHITE

# Data rows — color-coded cells
def _cell_color(v):
    """Heat-map: green for >1.20, lighter green for 1.05-1.20, neutral for 0.95-1.05, light red for <0.95."""
    if v >= 1.40:
        return RGBColor(0x00, 0x86, 0x3D)   # strong green
    if v >= 1.20:
        return RGBColor(0x6F, 0xC1, 0x80)   # green
    if v >= 1.05:
        return RGBColor(0xC9, 0xE7, 0xCB)   # light green
    if v >= 0.95:
        return RGBColor(0xF5, 0xF5, 0xF5)   # near-neutral
    if v >= 0.80:
        return RGBColor(0xFC, 0xE3, 0xCD)   # light orange
    return RGBColor(0xF4, 0xC2, 0xA1)       # orange (below 0.80)

for ri, (cat, payload) in enumerate(sorted_cats, start=1):
    p_arr = payload["profile"]
    s = payload["stats"]
    peak_m = months[s["peak_month"] - 1]
    is_override = bool(s.get("planner_override"))

    # Category cell
    cell = tbl.cell(ri, 0)
    tf = cell.text_frame; tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(1)
    tf.text = cat + ("  ★" if is_override else "")
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.runs[0]; r.font.size = Pt(8); r.font.bold = True; r.font.color.rgb = NAVY

    # Month cells
    for ci, v in enumerate(p_arr, start=1):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _cell_color(v)
        tf = cell.text_frame; tf.margin_left = tf.margin_right = Pt(1)
        tf.margin_top = tf.margin_bottom = Pt(1)
        tf.text = f"{v:.2f}"
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]; r.font.size = Pt(8); r.font.color.rgb = NAVY

    # #SKU cell
    cell = tbl.cell(ri, 13)
    tf = cell.text_frame; tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    tf.text = str(s["consistent_skus"])
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]; r.font.size = Pt(8); r.font.color.rgb = GREY

    # Peak cell
    cell = tbl.cell(ri, 14)
    tf = cell.text_frame; tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    tf.text = peak_m
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]; r.font.size = Pt(8); r.font.bold = True; r.font.color.rgb = ORANGE

# Legend
legend = slide.shapes.add_textbox(Inches(0.25), Inches(7.05), Inches(12.83), Inches(0.3))
legend.text_frame.text = (
    "★ = Planner override   ·   Heatmap: dark green ≥1.40 · green ≥1.20 · light green ≥1.05 · neutral ≥0.95 · light orange ≥0.80 · orange <0.80"
)
legend.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
legend.text_frame.paragraphs[0].runs[0].font.italic = True
legend.text_frame.paragraphs[0].runs[0].font.color.rgb = GREY


# ── Save ──
prs.save(str(DECK_PATH))
print(f"[ok] Appended 3 slides to {DECK_PATH} (now {len(prs.slides)} slides total)")
