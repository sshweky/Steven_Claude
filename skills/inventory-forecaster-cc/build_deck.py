"""
Builds the Inventory Forecaster — VP Review Deck (.pptx).
Source of truth = METHODOLOGY.md + recent back-test outputs.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Pets+People-ish palette (neutral pro)
NAVY    = RGBColor(0x1F, 0x3A, 0x5F)
ORANGE  = RGBColor(0xE8, 0x7A, 0x2C)
GRAY    = RGBColor(0x55, 0x5C, 0x66)
LIGHT   = RGBColor(0xF4, 0xF6, 0xF9)
GREEN   = RGBColor(0x2E, 0x8B, 0x57)
RED     = RGBColor(0xC0, 0x39, 0x2B)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x10, 0x14, 0x1A)


def add_title_bar(slide, prs, title, subtitle=None):
    # Top color band
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                   prs.slide_width, Inches(0.9))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    tf = band.text_frame
    tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.18)
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = WHITE
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.4), Inches(0.95),
                                        prs.slide_width - Inches(0.8), Inches(0.4))
        sub.text_frame.text = subtitle
        sp = sub.text_frame.paragraphs[0]
        sp.font.size = Pt(12); sp.font.italic = True; sp.font.color.rgb = GRAY


def add_footer(slide, prs, text):
    box = slide.shapes.add_textbox(Inches(0.3), prs.slide_height - Inches(0.4),
                                    prs.slide_width - Inches(0.6), Inches(0.3))
    tf = box.text_frame; tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(9); p.font.color.rgb = GRAY


def add_bullets(slide, prs, items, top=Inches(1.5), left=Inches(0.5),
                width=None, height=None, font_size=14, bullet_indent=0):
    if width is None:  width  = prs.slide_width  - Inches(1.0)
    if height is None: height = prs.slide_height - top - Inches(0.7)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # support tuples: (level:int, text, [color]) or (text, color)
        if isinstance(item, tuple):
            if len(item) == 1:
                level, text, color = 0, item[0], None
            elif isinstance(item[0], int):
                level = item[0]
                text  = item[1]
                color = item[2] if len(item) > 2 else None
            else:
                # treat as (text, color)
                level, text = 0, item[0]
                color = item[1] if len(item) > 1 else None
        else:
            level, text, color = 0, item, None
        if text == "":
            p.text = ""
        else:
            p.text = ("    " * level) + ("• " if level == 0 else "– ") + text
        p.font.size = Pt(font_size if level == 0 else max(11, font_size - 2))
        p.font.color.rgb = color or BLACK
        p.space_after = Pt(4)


def add_kpi_row(slide, prs, kpis, top=Inches(1.6)):
    """kpis: list of (label, value, color)."""
    n = len(kpis)
    margin = Inches(0.5)
    gap    = Inches(0.2)
    total_w = prs.slide_width - 2 * margin - gap * (n - 1)
    card_w = int(total_w / n)
    card_h = Inches(1.5)
    for i, (label, value, color) in enumerate(kpis):
        x = margin + i * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top,
                                       card_w, card_h)
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = color
        card.line.width = Pt(1.5)
        card.shadow.inherit = False
        # reset card text
        ctf = card.text_frame
        ctf.text = ""
        ctf.margin_top = Inches(0.15)
        ctf.margin_left = Inches(0.15)
        ctf.margin_right = Inches(0.15)
        # value (big)
        p = ctf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = value
        p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = color
        # label (small)
        p2 = ctf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = label
        p2.font.size = Pt(11); p2.font.color.rgb = GRAY


def add_table(slide, prs, headers, rows, top=Inches(1.5), left=Inches(0.5),
              width=None, font_size=11):
    if width is None:
        width = prs.slide_width - Inches(1.0)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    height = Inches(0.4 + 0.32 * n_rows)
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    for j, h in enumerate(headers):
        c = table.cell(0, j)
        c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(font_size + 1)
                r.font.color.rgb = WHITE
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            c = table.cell(i + 1, j)
            c.text = str(v)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(font_size)
                    r.font.color.rgb = BLACK
            if i % 2 == 0:
                c.fill.solid(); c.fill.fore_color.rgb = LIGHT


# ─────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)   # widescreen
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]

# ── Slide 1: Title ──────────────────────────────────────────────
s = prs.slides.add_slide(blank)
# big navy banner
b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
b.fill.solid(); b.fill.fore_color.rgb = NAVY; b.line.fill.background()
# orange accent bar
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.7),
                         Inches(0.8), Inches(0.08))
acc.fill.solid(); acc.fill.fore_color.rgb = ORANGE; acc.line.fill.background()

t1 = s.shapes.add_textbox(Inches(0.6), Inches(1.5), prs.slide_width - Inches(1.2), Inches(1.2))
tf = t1.text_frame; tf.text = "Inventory Forecaster"
tf.paragraphs[0].font.size = Pt(54); tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = WHITE

t2 = s.shapes.add_textbox(Inches(0.6), Inches(2.85), prs.slide_width - Inches(1.2), Inches(0.6))
tf = t2.text_frame
tf.text = "VP-of-Planning Feedback Stack — Q1 / Q2 / Q3 / Q4"
tf.paragraphs[0].font.size = Pt(24); tf.paragraphs[0].font.color.rgb = ORANGE

t3 = s.shapes.add_textbox(Inches(0.6), Inches(3.6), prs.slide_width - Inches(1.2), Inches(2.5))
tf = t3.text_frame; tf.word_wrap = True
for line in [
    "Pets+People  ·  Acct 1864 (Amazon)  ·  May 2026",
    "",
    "What changed, why, and how each change behaved on a 1,511-record back-test.",
]:
    p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
    p.text = line
    p.font.size = Pt(16); p.font.color.rgb = WHITE

# ── Slide 2: Why we changed anything ────────────────────────────
s = prs.slides.add_slide(blank)
add_title_bar(s, prs, "Why we changed anything",
              "Four targeted VP-of-Planning comments on the production forecaster")
add_bullets(s, prs, [
    "VP-Q1: \"Baseline-mode logic is biased low — recent OOS quiet weeks are dragging the average.\"",
    "VP-Q2: \"OOS cancellations look like zero demand. Forecaster should reconstruct true intent.\"",
    "VP-Q3: \"Bi-weekly cadence enforcement isn't necessary — only enforce monthly+ patterns.\"",
    "VP-Q4: \"Don't double-count confirmed customer POs — zero out AI in any week with an open PO.\"",
    "",
    "Each became a code change on `inventory_forecaster.py` and a back-test on acct 1864.",
    (1, "Goal: deliver four small, defensible improvements that planners can audit on a single record."),
], top=Inches(1.4), font_size=16)
add_footer(s, prs, "inventory-forecaster-cc · METHODOLOGY.md")

# ── Slide 3: Pipeline at a glance ───────────────────────────────
s = prs.slides.add_slide(blank)
add_title_bar(s, prs, "Pipeline at a glance",
              "What runs end-to-end on every record")
add_bullets(s, prs, [
    "Pull projections + Order_History + Amazon_Catalog from Quickbase via CData",
    "Classify into one of 9 model types (Croston's, Seasonal Baseline, HW, Sparse, Heuristic, …)",
    "Run the matched model on a weighted 78-obs history (52w + L13W appended 2× for 3× weight)",
    "Apply event lifts: Prime Day W7-W9 (+25%, Amazon only) · Fall Deal W23-W25 (+12%)",
    "Snap each non-zero week to the master-pack multiple",
    (1, "VP-Q4: zero-out forward weeks with confirmed customer POs", ORANGE),
    "Compare AI vs manual projections → AI_ALERT if >5% variance",
    "Write back AI_PRJ_W1..W26 + AI_ALERT (parallel UPDATE, resume-safe)",
], top=Inches(1.4), font_size=15)
add_footer(s, prs, "Phases 1–4 of inventory_forecaster.py")

# ── Slide 4: Model classification ───────────────────────────────
s = prs.slides.add_slide(blank)
add_title_bar(s, prs, "Model classification — deterministic routing",
              "Evaluated in order. First match wins.")
add_table(s, prs, ["Test", "Model", "Typical pattern"], [
    ["Zero L13W AND L26W",            "Inactive",            "Discontinued / out of policy"],
    ["≤6 weeks of history",           "New / Relaunch",      "Just hitting the catalog"],
    ["Recently inactive, restarting", "Reactivating",        "Returning after pause"],
    ["<13 active weeks in L52W",      "Sparse Intermittent", "Long tail / accessory"],
    ["Dense (≥50% non-zero L13W)",    "Seasonal Baseline",   "Smooth, steady mover"],
    ["CV>0.5 OR zeros>20%",           "Croston's",           "Lumpy / intermittent"],
    ["CV≤0.5 AND zeros≤20%",          "Holt-Winters",        "Steady with trend"],
    ["Out-of-policy / EOL",           "OTB (zero)",          "Buy-out / clearance"],
], top=Inches(1.5), font_size=12)
add_footer(s, prs, "9 models, deterministic — no hidden ML")

# ── Slide 5: VP-Q1 ──────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_title_bar(s, prs, "VP-Q1 — Evidence-based baseline mode",
              "L13W non-zero average instead of all-weeks average")
add_bullets(s, prs, [
    ("BEFORE: baseline = L13W all-weeks average."),
    (1, "Post-Prime Day quiet weeks (legitimate drawdown, not lost demand) dragged the mean down 15-20%.", GRAY),
    "",
    ("AFTER: baseline = L13W non-zero average. Fallback chain: L26W non-zero → L13W all-weeks → last-resort."),
    (1, "Reflects true per-order rate. Quiet weeks no longer suppress the baseline.", GREEN),
    "",
    "Back-test impact (acct 1864):",
    (1, "+18.5% aggregate 26w demand · 1,199 records lifted · broadly distributed across all customers."),
    (1, "Largest individual lifts on Amazon items recovering from drawdown — exactly the targeted population."),
], top=Inches(1.4), font_size=15)
add_footer(s, prs, "scripts/inventory_forecaster.py · seasonal_baseline()")

# ── Slide 6: VP-Q2 ──────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_title_bar(s, prs, "VP-Q2 — OOS-aware demand reconstruction",
              "Pull Order_History per-week; classify cancels by reason code")
add_bullets(s, prs, [
    "Bucket A — OOS-driven (Inventory Error, Supplier Delay, Slot Constraint…):",
    (1, "KEEP demand intent. The customer wanted units — we just didn't ship.", GREEN),
    "Bucket B — Demand-invalidating (Customer Order Error, Future Delete, Low Margin):",
    (1, "SUBTRACT from clean demand. The order was never real intent.", RED),
    "Bucket C — Ambiguous (Other, Any, null):",
    (1, "Keep as-is. Conservative.", GRAY),
    "",
    "`clean_ord` replaces `raw_ord` in the model fit. OOS severity ≥15% in any week is logged as a driver.",
    (1, "Per-record swings can be large for chronic-OOS items; aggregate impact is modest (≤2% on acct 1864)."),
], top=Inches(1.4), font_size=14)
add_footer(s, prs, "scripts/oos_history.py · fetch_clean_demand() + classify_cancel()")

# ── Slide 7: VP-Q3 ──────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_title_bar(s, prs, "VP-Q3 — Cadence relaxation",
              "Stop synthesizing zeros for items that ordered weekly with one quiet parity")
add_bullets(s, prs, [
    "BEFORE: any record with ≥70% zeros on one parity over L26W was flagged \"bi-weekly\" and forced into a zero/order/zero pattern.",
    (1, "Over-constrained: a recent recovery from 6 quiet weeks looked exactly like bi-weekly cadence.", RED),
    "",
    "AFTER: detect_biweekly() only fires for monthly+ cadences (median gap ≥3 weeks AND ≥60% gap consistency).",
    (1, "Returns the median gap (3, 4, 5, …). apply_ordering_pattern() does N-week chunk merging.", GREEN),
    "",
    "Back-test impact:",
    (1, "0 records flagged as \"bi-weekly\" (down from ~88 under old logic on acct 1864)."),
    (1, "Records with monthly cadence still get appropriate gap-aware merging."),
], top=Inches(1.4), font_size=15)
add_footer(s, prs, "detect_biweekly() · apply_ordering_pattern()")

# ── Slide 8: VP-Q4 ──────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_title_bar(s, prs, "VP-Q4 — Don't double-count confirmed customer POs",
              "Strict zero-out in any forward week with an open PO")
add_bullets(s, prs, [
    "Pull every Order_History row with `Qty_Open > 0` per Acct_MStyle.",
    "Bucket each by **Cancel_Date** = customer ship-by deadline.",
    (1, "Critical: NOT Next_Rcpt_Date — that's incoming supplier merch, wrong direction.", RED),
    "",
    "For any forward week 1-26 where the bucketed open-PO qty > 0:",
    (1, "fcst[week] = 0 (strict zero, not subtract).", ORANGE),
    "",
    "Why strict zero, not subtract:",
    (1, "The confirmed PO IS the demand signal for that week. Replen counts the PO already.", GRAY),
    (1, "AI projection on top would systematically over-buy front weeks.", GRAY),
    "",
    "On acct 1864: 351 keys had confirmed forward POs · 134,888 total open units across 1,937 PO lines.",
], top=Inches(1.4), font_size=14)
add_footer(s, prs, "scripts/oos_history.py · fetch_open_pos_forward()")

# ── Slide 9: Back-test headlines ────────────────────────────────
s = prs.slides.add_slide(blank)
add_title_bar(s, prs, "Back-test headlines — acct 1864 (1,511 records)",
              "Aggregate 26-week demand impact, isolating each VP comment")
add_kpi_row(s, prs, [
    ("VP-Q1 baseline lift",  "+18.5%",  GREEN),
    ("VP-Q2 OOS reconstruct", "≤2%",    GRAY),
    ("VP-Q3 cadence cleanup", "≈ flat", GRAY),
    ("VP-Q4 PO units found",  "134,888", ORANGE),
])
add_bullets(s, prs, [
    "VP-Q1 is the dominant aggregate move. Q2/Q3/Q4 produce per-record corrections, not headline shifts.",
    "VP-Q2 swings are concentrated on chronic-OOS Amazon SKUs (BB35096, FF6636AMZ, FF8424).",
    "VP-Q4 affects 351 of 1,511 records (23%) — front-week zeroing where replen would otherwise double-count.",
    "",
    ("CAVEAT: latest two back-test runs (v7, v8, v9) hit transient CData IncompleteRead failures on the second", RED),
    ("subprocess. Production guards now abort instead of silently producing contaminated forecasts.", RED),
], top=Inches(3.4), font_size=13)
add_footer(s, prs, "backtest_vp_q1_acct1864.md · backtest_vp_q2_acct1864.md")

# ── Slide 10: Reliability hardening ─────────────────────────────
s = prs.slides.add_slide(blank)
add_title_bar(s, prs, "Reliability hardening (May 2026)",
              "Caught a silent-failure pattern before write-back")
add_bullets(s, prs, [
    "What we found: under transient CData IncompleteRead failure, the forecaster used to continue with empty",
    (1, "oos_data {} and open_pos_data {} dicts — silently producing forecasts that skipped Q2 and Q4 entirely.", RED),
    "",
    "Fix 1 — abort-on-empty guards:",
    (1, "If --oos-smoothing requested and oos_data is empty → sys.exit() with [ABORT] message.", GREEN),
    (1, "If VP-Q4 enabled (default) and open_pos_data is empty → sys.exit() too. Use --no-po-zero to opt out.", GREEN),
    "",
    "Fix 2 — better retry budget:",
    (1, "MAX_RETRIES: 3 → 5 (62s total backoff vs 14s).", GREEN),
    (1, "On IncompleteRead, force re-prime of the CData session before retry — old code reused a dead socket.", GREEN),
    "",
    "Net effect: a contaminated run can no longer reach write-back.",
], top=Inches(1.4), font_size=14)
add_footer(s, prs, "cdata_query() · Phase 2.7 + 2.8 guards")

# ── Slide 11: What we recommend ─────────────────────────────────
s = prs.slides.add_slide(blank)
add_title_bar(s, prs, "Recommended next steps",
              "Path from back-test to production write-back")
add_bullets(s, prs, [
    "Wait for CData to stabilize (Phase 1 has been intermittent today).",
    "Run a clean v10 back-test once both subprocess runs complete without [FAIL] lines.",
    "Spot-check 3–5 records in the viewer (Croston's, Seasonal Baseline, Sparse Intermittent).",
    "",
    "Production write-back command:",
    (1, "python scripts/inventory_forecaster.py --acct 1864 --oos-smoothing"),
    (1, "(VP-Q1 + Q3 are baked in; --oos-smoothing turns on Q2; Q4 is on by default.)"),
    "",
    "Roll-out:",
    (1, "Acct 1864 first (current scope) → review for 1 week → expand to other accounts."),
    (1, "AI_ALERT >5% surfaces variance for planner review without blocking write-back."),
], top=Inches(1.4), font_size=15)
add_footer(s, prs, "Production-ready after one clean back-test run")

# ── Slide 12: Appendix — config knobs ───────────────────────────
s = prs.slides.add_slide(blank)
add_title_bar(s, prs, "Appendix — configuration knobs",
              "Constants in scripts/inventory_forecaster.py")
add_table(s, prs, ["Constant", "Value", "Purpose"], [
    ["PRIME_DAY_WEEKS",    "{7, 8, 9}",   "Amazon-only pre-order window"],
    ["FALL_DEAL_WEEKS",    "{23, 24, 25}", "All-account pre-order window"],
    ["PRIME_DAY_LIFT",     "1.25",        "+25% on Prime Day weeks"],
    ["FALL_DEAL_LIFT",     "1.12",        "+12% on Fall Deal weeks"],
    ["AMAZON_CUST_SUBSTR", "\"AMAZON\"",  "Gates Q2-Amazon POS pulls + Prime Day lifts"],
    ["DAMP",               "0.1",         "Seasonal profile dampening (±20% from 1.0)"],
    ["ALERT_THRESHOLD",    "0.05",        "5% variance triggers AI_ALERT"],
    ["MAX_RETRIES",        "5",           "CData retry budget (was 3 before May 2026)"],
], top=Inches(1.5), font_size=12)
add_footer(s, prs, "Appendix · METHODOLOGY.md §9")


out = r"C:\Users\steven\.claude\skills\inventory-forecaster-cc\inventory_forecaster_VP_review.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
