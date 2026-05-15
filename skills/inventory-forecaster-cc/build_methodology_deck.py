#!/usr/bin/env python3
"""
build_methodology_deck.py
-------------------------
Generates inventory_forecaster_methodology.pptx — a plain-English walk-through
of the inventory forecaster for novice planners.  No jargon, no math notation.
Every refinement is explained in business terms.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = r"C:\Users\steven\.claude\skills\inventory-forecaster-cc\inventory_forecaster_methodology.pptx"

# ── Brand palette ─────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1F, 0x3A, 0x5F)
TEAL      = RGBColor(0x2A, 0x9D, 0x8F)
ORANGE    = RGBColor(0xE7, 0x6F, 0x51)
GOLD      = RGBColor(0xE9, 0xC4, 0x6A)
GREY_DK   = RGBColor(0x33, 0x33, 0x33)
GREY_MD   = RGBColor(0x66, 0x66, 0x66)
GREY_LT   = RGBColor(0xF2, 0xF2, 0xF2)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
RED       = RGBColor(0xC6, 0x28, 0x28)
GREEN     = RGBColor(0x2E, 0x7D, 0x32)

# ── Slide constants (16:9) ────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.5)
CONTENT_W = SLIDE_W - 2 * MARGIN

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill=NAVY, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    return sh

def add_text(slide, left, top, width, height,
             text, *, size=14, bold=False, color=GREY_DK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb

def add_bullets(slide, left, top, width, height, items, *,
                size=12, color=GREY_DK, bullet_color=TEAL):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(5)
        bul = p.add_run()
        bul.text = "●  "
        bul.font.name = "Calibri"
        bul.font.size = Pt(size)
        bul.font.color.rgb = bullet_color
        bul.font.bold = True
        run = p.add_run()
        run.text = item
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb

def slide_header(slide, title, subtitle=None, *, color=NAVY):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), fill=color)
    add_text(slide, MARGIN, Inches(0.10), CONTENT_W, Inches(0.45),
             title, size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, MARGIN, Inches(0.50), CONTENT_W, Inches(0.30),
                 subtitle, size=11, color=GOLD, anchor=MSO_ANCHOR.TOP)
    add_rect(slide, 0, Inches(7.30), SLIDE_W, Inches(0.04), fill=TEAL)

def page_footer(slide, n, total):
    add_text(slide, Inches(12.6), Inches(7.10), Inches(0.7), Inches(0.25),
             f"{n} / {total}", size=9, color=GREY_MD, align=PP_ALIGN.RIGHT)
    add_text(slide, MARGIN, Inches(7.10), Inches(8.0), Inches(0.25),
             "Inventory Forecaster · Plain-English Guide for Planners",
             size=9, color=GREY_MD)

def add_table(slide, left, top, width, height, headers, rows, *,
              header_fill=NAVY, header_font=WHITE, alt_fill=GREY_LT,
              size=10, hdr_size=11):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers),
                                         left, top, width, height).table
    for j, h in enumerate(headers):
        cell = table_shape.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        tf = cell.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = h
        r.font.name = "Calibri"; r.font.bold = True
        r.font.size = Pt(hdr_size); r.font.color.rgb = header_font
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table_shape.cell(i + 1, j)
            if i % 2 == 1:
                cell.fill.solid(); cell.fill.fore_color.rgb = alt_fill
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame; tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            r.font.name = "Calibri"; r.font.size = Pt(size)
            r.font.color.rgb = GREY_DK
    return table_shape

# ── Refinement-card helper ────────────────────────────────────────────────────
# Layout v4 (2026-05-06): tightened section spacing; EXAMPLE rendered INSIDE
# the card box with its own ORANGE label; acct-mstyle key bolded NAVY for
# scanability.  Designed for h=2.10" cards.
def refinement_card(slide, left, top, w, h, code, title, before, after, why, example=""):
    add_rect(slide, left, top, w, h, fill=WHITE, line=GREY_LT)
    # Left ribbon w/ code
    add_rect(slide, left, top, Inches(0.85), h, fill=NAVY)
    add_text(slide, left, top, Inches(0.85), h, code, size=14, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Title strip (compact)
    add_text(slide, left + Inches(0.95), top + Inches(0.02), w - Inches(1.0), Inches(0.22),
             title, size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.TOP)
    # Section geometry
    SEC_LBL_W   = Inches(0.95)
    SEC_OFF     = Inches(1.00)
    SEC_TXT_OFF = Inches(2.00)
    SEC_TXT_W   = w - Inches(2.05)
    # Vertical layout — sections packed tight:
    #   title 0.26 + problem 0.34 + fix 0.55 + why 0.36 + example 0.59 = 2.10"
    y = top + Inches(0.26)
    # PROBLEM
    add_text(slide, left + SEC_OFF, y, SEC_LBL_W, Inches(0.18),
             "PROBLEM", size=8, bold=True, color=RED)
    add_text(slide, left + SEC_TXT_OFF, y, SEC_TXT_W, Inches(0.34),
             before, size=10, color=GREY_DK)
    y += Inches(0.34)
    # FIX
    add_text(slide, left + SEC_OFF, y, SEC_LBL_W, Inches(0.18),
             "FIX", size=8, bold=True, color=GREEN)
    add_text(slide, left + SEC_TXT_OFF, y, SEC_TXT_W, Inches(0.55),
             after, size=10, color=GREY_DK)
    y += Inches(0.55)
    # WHY YOU CARE
    add_text(slide, left + SEC_OFF, y, SEC_LBL_W, Inches(0.18),
             "WHY YOU CARE", size=8, bold=True, color=TEAL)
    add_text(slide, left + SEC_TXT_OFF, y, SEC_TXT_W, Inches(0.36),
             why, size=10, color=GREY_DK)
    y += Inches(0.36)
    # EXAMPLE — inside the card, with key bolded
    if example:
        add_text(slide, left + SEC_OFF, y, SEC_LBL_W, Inches(0.18),
                 "EXAMPLE", size=8, bold=True, color=ORANGE)
        # Remaining card height for example body
        ex_h = (top + h) - y - Inches(0.04)
        if ex_h < Inches(0.20):
            ex_h = Inches(0.20)
        ex_tb = slide.shapes.add_textbox(left + SEC_TXT_OFF, y, SEC_TXT_W, ex_h)
        tf = ex_tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
        tf.margin_top  = Inches(0.0);  tf.margin_bottom = Inches(0.02)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        # Convention: example is "Acct N, Mstyle X — long detailed explanation"
        # Bold the "Acct N, Mstyle X" prefix in NAVY for instant scannability.
        if " — " in example:
            key_part, rest = example.split(" — ", 1)
            sep = " — "
        elif " - " in example:
            key_part, rest = example.split(" - ", 1)
            sep = " — "
        else:
            key_part, rest, sep = example, "", ""
        r1 = p.add_run(); r1.text = key_part
        r1.font.name = "Calibri"; r1.font.size = Pt(10)
        r1.font.bold = True; r1.font.color.rgb = NAVY
        if rest:
            r2 = p.add_run(); r2.text = sep + rest
            r2.font.name = "Calibri"; r2.font.size = Pt(10)
            r2.font.color.rgb = GREY_DK

# ── Example acct-mstyle keys per rule ─────────────────────────────────────────
# Real keys harvested from a dry-run forecast against Amazon + Walmart +
# Burlington + Petbarn (2026-05-06).  Each entry is one key per rule that
# actually fires that rule in production data.  When a rule has no harvested
# example yet, the EXAMPLES.get(...) default ("") suppresses the line.
EXAMPLES = {
    # Real keys harvested from instrumented dry-run forecasts (2026-05-06) across
    # AMAZON / WAL MART / BURLINGTON / CHEWY (2,085 records).  Each entry comes
    # from the per-record `rule_fires` list with actual harvest numbers cited.
    # Format: "Acct N, Mstyle X — detailed plain-English explanation"
    # The "Acct N, Mstyle X" prefix renders bold-NAVY, the rest GREY_DK.
    "VP-Q1":   "Acct 23011, Mstyle BB0237 — Walmart Seasonal Baseline. Baseline mode 'L13 nz-avg (OOS: fill-rate 8% over 9,024 units)'. The all-weeks average would have been dragged down by post-event quiet weeks where Walmart drew down stock; we use the non-zero average (5,760 forecast units vs 5,150 manual flat-200/wk plan) so the baseline reflects true per-order rate.",
    "VP-Q2":   "Acct 23011, Mstyle FF12660 — Walmart, 3% fill-rate over 45,720 units (massive backlog from a long stockout). VP-Q2 keeps the unfilled order qty as 'demand intent' rather than reading the missed weeks as zero demand. AI projects 75,036 units vs 40,258 planned (+86%) because the per-order true rate is ~2x the suppressed all-weeks average.",
    "VP-Q3":   "Acct 23011, Mstyle BB33708 — Walmart Croston's, fired with a drawdown signature: ≥3-week consistent gap detected, so cadence enforcement kicks in. Without VP-Q3 the system would have inserted phantom zero weeks; here the rule confirms the gap is real (post-event drawdown) and smooths around it instead of forcing alternation.",
    "VP-Q4":   "Acct 23011, Mstyle BB13435CLR/12 — Walmart Croston's. There is a confirmed open PO landing in W3, so VP-Q4 zeroes out that specific week's AI projection (the PO already covers the demand). Without this rule, replenishment would double-count the PO + AI projection and overstock the DC by 1 week's worth.",
    "R1":      "Acct 13640, Mstyle FF38675 — Burlington Coat Factory, model 'OTB (zero)'. Classic off-price pattern: 1-2 big lump orders in L52, then dormant for 4+ months. Forecast = 0 for 26w because Burlington buys in occasional bulk closeouts, not weekly cadence. Manual had 7,800 units that R1 wiped out — Burlington isn't going to reorder this week-by-week. Amazon gate (2026-05-07): Amazon items are NEVER classified as OTB regardless of pattern — Amazon ordering is centrally managed and even sparse-looking histories reflect ongoing replenishment. Amazon items that look 'OTB-shaped' route to the standard Inactive recipe instead, which restarts forecasting the moment orders resume.",
    "R2":      "Acct 23011, Mstyle BB0092PCS2 — Walmart Sparse Intermittent. R2's 1.5× sanity ceiling binds: a single big historical order would otherwise inflate the 26-week total to ~38% above what the account's average pace can support over 6 months. Cap pulls it back to a defensible level — no phantom amplification from one outlier order.",
    "R3":      "Acct 1864, Mstyle BB35237 — model 'Inactive+Floor (R3)'. Customer ordered ≥4 times in last 6 months AND ≥8 times in last year, so instead of forcing AI=0, R3 applies a small steady floor. Catches the slow-but-real seller that VP-Q1's strict Inactive rule would otherwise zero out. Avoids leaving a real account dry.",
    "R5":      "Any international retailer (e.g. Petbarn AU, Loblaws CA, Comercializadora MX). They place big orders quarterly — at week 14 of silence the standard 13-week Inactive rule would mark them dead. R5 bumps the threshold to 26 weeks, so a customer who naturally orders every 14-20 weeks still gets a forecast through the quiet quarter.",
    "R6":      "Acct 23011, Mstyle FF31287 — Walmart Croston's, ordered 7 of last 13 weeks at ~13 units/order. R6's high-volume override boosts the predicted order size toward the recent average instead of letting Croston's z-estimate drift down. AI lands at 208 units vs the manual flat 23/wk plan (which would have over-projected for an account in the middle of a draw-down).",
    "R7":      "Any Amazon item with zero historical signal in W23-W25 (Fall Deal weeks). Without history, the AI would skip the Fall Deal entirely. R7 inserts a synthetic order in the W23-W25 window scaled at +12% lift over the L13 baseline, so Amazon items still participate in the second-biggest event of the year even when last year's history is missing.",
    "R8":      "A record with ≥9 active weeks where the top 2 orders make up >60% of L13 total (e.g. one 800-unit order, one 600-unit order, the rest 50-100 units). The plain average gets hijacked by the two bursts. R8 swaps to median per-order, which tells the truth about the repeat-order size and ignores the two outliers.",
    "R9":      "Acct 3102, Mstyle FF33495 — Chewy Heuristic. Reactivating-recipe item bouncing back from a quiet period. R9 caps the baseline at 2× the L52W weekly average so the recovering forecast can't run hotter than 2× its long-run pace. AI lands at 8 units vs 558 manual (Chewy's flat 30/wk plan was way over the actual L13 sparse rate of 4/wk). 2026-05-07 fix: R9 is now applied UNCONDITIONALLY (was previously skipped when F23b trailing-zero discount also fired, which let single-PO patterns escape the ceiling). Concrete case: FF7612 Petco had L13 = single 5,208-unit PO ~12w ago then dormant; F23b discounted baseline by 0.30× to 1,562, which projected ~37K vs L26 all-weeks avg 426/wk. Multiplier raised to 2.5× when F23b also fires, to soften the double-discount.",
    "F4":      "Acct 1864, Mstyle FF31061 — Amazon, baseline mode 'L26 nz-avg (sparse L13)'. L13 had fewer than 5 active weeks but L52 had 8+ active weeks. F4 detected the thin recent quarter and pulled in the 26-week non-zero average as the floor — so a quiet quarter doesn't make a real product look dead. AI 534 units vs 990 manual flat 40/wk.",
    "F6 / F26 / F27": "Acct 23011, Mstyle BB0237 — Walmart Seasonal Baseline. F6 detected L4 non-zero average ≤ 50% of L13 non-zero average — a real recent slowdown — and multiplied the baseline by 0.65. (F26: 50-70% threshold → 0.85x. F27: reverse for ramp-up where L4 is 30-60% above L13 → 1.10x.) Forecast leans into the last month, not just the last quarter.",
    "F25":     "A record with one big freak order (e.g. 744 units) when normal weeks are ~50/wk and there are 4+ other healthy weeks. F25 drops the outlier from the L13 average entirely (rather than just capping it). One weird buy doesn't get to set the next 6 months of forecast.",
    "F24":     "Acct 23011, Mstyle BB0096PCS2 — Walmart Seasonal Baseline, fill-rate 7% over 408 units. After POS blend, drawdown lift, and order-coverage lift stacked, the baseline drifted to ~3-4× the recent L13 rate. F24 caps final baseline at L13 avg × 2.0 — flat history can't be inflated indefinitely. AI 1,404 units vs 1,040 manual flat 40/wk (+35%).",
    "F7":      "A grill brush in week 5 (winter trough, ~10/wk) with a known summer peak of ~30/wk. F7 detects the seasonal category AND the steep peak/trough ratio, so instead of anchoring on the dead winter rate it anchors on the historical peak adjusted by current week's seasonal factor. Item entering peak gets a realistic forecast, not one stuck on the trough rate.",
    "F11":     "Any Amazon item in W5-W9 (Prime Day pre-buy weeks). Instead of a flat +25% lift across W7-W9, F11 applies a ramp: W5 +10%, W6 +15%, W7 +25%, W8 +25%, W9 +20%. Matches what planners actually see — orders ramp up toward peak then taper, not a step function.",
    "F16 / F16b": "Acct 23011, Mstyle BB0237 — steep peak/trough ratio AND L13 avg ≥ 50/wk AND seasonal category match. F16 disables the standard ±20% damping and lets the raw seasonal shape pass through with the relaxed 0.30-2.5× cap, so the natural 3× summer peak isn't crushed to a flat profile.",
    "F13":     "Acct 1864, Mstyle BB22116 — Amazon Seasonal Baseline. POS sell-through rate is meaningfully higher than the order rate AND POS is healthy — Amazon is burning down inventory faster than they're ordering from us. F13 raises the baseline to anticipate the impending refill order. Catches Amazon's classic 'drawdown then big refill' pattern.",
    "F14a / F14b": "Acct 1864, Mstyle BB0012AMZ6 — Amazon Croston's. POS L4/L13 ratio = 1.0 (consumers still buying). F14a blocks F10's decline-detection from scaling the item down on a buyer-side ordering pause: trust consumer sales over a buyer-side hiccup. Item doesn't get prematurely killed by a 1-2 week ordering lull when POS proves it's still alive.",
    "F15":     "An Amazon item with order rate consistently >1.15× the consumer POS rate — buyer is keeping a structural premium for safety stock or distribution coverage. F15 blends the baseline toward the higher order rate so the AI respects the buyer's choice rather than chasing only consumer demand.",
    "F18":     "An Amazon Lumpy-recipe item with healthy POS and consumer sales running faster than the order pattern would predict. F18 raises the predicted order size to match consumer velocity. Forecast catches up to actual consumer sell-through instead of staying anchored on the (lagging) order pattern.",
    "T4":      "A Chewy/Petco.com/PetSmart.com Seasonal Baseline item where L4 non-zero avg is running ≥5% above L13. These ecommerce stores have no consumer-sales feed, but planners see late-cycle ramps. T4 blends the baseline toward the recent rate (heavier weight if L4 is ≥15% above L13), capturing acceleration the AI can't see from raw orders.",
    "F10":     "Acct 23011, Mstyle BB13435CLR/12 — Walmart Croston's. F10's two checks both passed: L4 ≤ 70% of L13 AND year-over-year is also down. Without both checks, a seasonal trough could trigger a false-positive decline; the YoY guard ensures we only scale down on a real decline. AI lands at 87,588 units vs 75,075 manual.",
    "F22a / F22c": "A record ending the recent history with 4-5 zero weeks in a row (mid-drawdown). F22a discounts the baseline by trailing-zero-count / 13. F22c additionally caps the final baseline at the recent average when the L13 window is thin (< 7 active weeks). Drawdown items get realistic next-quarter forecasts, not full-pace ones.",
    "F23a / F23b": "Acct 23011, Mstyle BB0094PCS2 — Walmart Heuristic on sparse history. F23b applied the same DAMP=0.85 seasonal shape used by Seasonal Baseline plus the trailing-zero discount (up to 70%), so a single big historical order doesn't get amplified 10× by seasonal positioning. AI 1,500 units vs 1,240 manual flat 50/wk (+21%).",
    "F9":      "A 15K+ annual-units sparse item where the L13 average is below both the L26 and L52 averages — recent quarter happened to be quiet. F9 picks the MAX of the three windows as the baseline so a strong-but-lumpy item doesn't get under-forecast by one quiet quarter that's not representative of true demand.",
    "F17 / F17b": "Acct 23011, Mstyle BB0100PCS2 — Walmart Sparse Intermittent with 25+/wk recent average. F17 shifted the next predicted order into W1 because the planner expected an order then. Low-volume tail items (under 25/wk) skip this — F17 only fires when there's enough real volume to justify pulling forward.",
    "F28":     "A Lumpy-recipe item where the model output would run 15%+ below the L13 quarterly average. F28 floors at a fixed share of L13 avg so Lumpy doesn't systematically under-forecast against the recent quarter pace. Closes a known under-forecast bias of ~511K units across 414 records pre-fix.",
    "F29":     "A new item that just started shipping with only 1 active week in last 4. The original rule needed 2 active weeks to apply the new-item floor. F29 loosens to use any shipped week from L4 or L8 — catches first-quarter shipping reality without over-restricting brand-new items.",
    "F30":     "Top-volume Seasonal-Baseline records (the 80th-percentile big items). F30 enforces a tighter baseline-vs-recent-rate ratio cap so big items can't run hot from stacked lifts. These are buyer-driven — the AI shouldn't out-forecast a human plan when the human has high-fidelity context.",
    "F30 (rev)": "Acct 23011, Mstyle BB0083 — model 'Inactive (zero order history)' with manual still showing 442 units. F30(rev) hard rule: zero orders in 26+ weeks → AI = 0, no exceptions. Doesn't matter if POS sells through — the customer hasn't ordered for half a year, POS alone isn't enough to project orders.",
    "F31":     "Acct 23011, Mstyle BB24711 — Walmart, model 'Pre-launch NEW (manual passthrough)'. PT_Item_Status flags this as new/launching, and there's zero order history. F31 passes the planner's manual through unchanged (75 units) instead of producing a bogus AI projection from no signal. Pre-launch items don't get random predictions.",
    "F32":     "A sparse item with near-zero recent history. F32 replaces the old sum-of-26-weeks cap (which rarely fired because sparse 26w totals are tiny) with a per-week clamp + a tiny-signal floor. Tighter, more reliable cap on small sparse items so a single big historical week doesn't blow up the forecast.",
    "F34":     "Acct 23011, Mstyle BB13435CLR/12 — Walmart Croston's. Weeks 27-51 ago sum to <1% of last 26 weeks (item launched ~26w ago). F34 marks it as a NEW LAUNCH and skips F10 decline-detection + the L52 ceiling that would have read those pre-launch zeros as a dying item and squashed the forecast.",
    "F35":     "Acct 23011, Mstyle FF12660 — Walmart Seasonal Baseline. A multi-week stockout gap was detected; F35 used the planner decay schedule (W1 stockout: 25% lost / 75% recoverable; W2: 50/50; W3: 75/25; W4+: 100% lost) to strip the recoverable-backlog portion of the catch-up week. L13 now reflects true per-order demand intent, not a backlog spike.",
    "F36":     "Acct 1864, Mstyle SF8169 — Amazon. ~28k units shipped 11-12 weeks ago against ~250/wk POS sell-through = ~112w of cover. F36 (Amazon-only) computes weeks-of-supply from a recent big shipment cluster, subtracts weeks elapsed, and forces AI W1..W{remaining} to 0 until consumer demand burns through. 2026-05-07 fixes: (1) SHP window widened from L13W→L26W so older big stockups are still detected (SF8169's stockup at W-11 was outside the prior 13w window). (2) Added L4 active-orders guard — if recent ORDERS run at ≥70% of POS rate the customer is actively replenishing, so F36 doesn't fire even if a recent shipment was big (catches cases like FF12853 where a big shipment landed the same week as ramping orders L4=2,600/wk vs POS=1,700/wk).",
    "F37":     "Acct 23011, Mstyle BB0096PCS2 — Walmart, anticipated on-hand (Inv_Wk6) goes negative. F37 caps that week's AI projection to what we can actually ship and rolls 75% of the unmet portion into W7's demand cohort (cohort decays 25% per week, fully lost at age 4w+). Forecast becomes physically realistic given inventory plan.",
    "F38":     "Acct 1864, Mstyle BB0150 — Amazon Seasonal Baseline, F38b fired. POS L4w +24% above L13w with buybox stable above MAP, so baseline lifted full-pct: AI 20,364 vs 16,435 manual (+24%). On the dry-run, F38 fired 236× across 1,513 Amazon records: F38b=102 (lifts), F38e=50 (cuts), F38f=64 (offline recovery), F38c=15, F38d=3, F38a=2.",
    "F39":     "Acct 1864, Mstyle FF7618 — Amazon Heuristic. F39 detected a duplicate-order run: the same large qty appeared in 3+ adjacent weeks within ±5%, flagging a phantom re-broadcast (one real order replayed by upstream feeds). The duplicate copies were stripped from history before model classification, so the L13/L26 baselines reflect actual demand, not feed echo.",
    "F40":     "Acct 1864, Mstyle FF7618 — order rate decelerating. After F39 dedup, the last 3 non-zero orders averaged 60/120/360 — a clear ramp-down pattern. F40 detected the deceleration and scaled the forward forecast to inherit the slowing pace, instead of letting Heuristic propagate the older heavier rate.",
    "F41":     "Acct 1864, Mstyle FF7618 — Amazon ship-lag dedupe. Shipment record for week LW_16 (14,328 units) effectively reappeared as a near-duplicate order at LW_15 (14,184 units, ~1% drift). F41 uses shipment evidence (the strongest signal) to keep LW_16 and zero LW_15 — phantom dedupe with ±15% tolerance. Runs BEFORE F39. Brought FF7618 from 86,580/26w → 17,760/26w.",
    "F42":     "Acct 1864, Mstyle SF8169 — Amazon Heuristic. After F41/F39 stripped phantoms, the 26w avg was still inflated by surviving lumpy orders. F42 anchors the Heuristic baseline to POS sell-through: ~250/wk × 26 = ~6,500 — vs the pre-cap Heuristic projection of 87,912. Amazon-only, only when POS data is reliable and coverage rules (F36/F38) aren't already governing zeroing.",
    "F43":     "Acct 23011, Mstyle FF25895 — Walmart. The last 4 weeks held a single 5×-median outlier that wasn't a real demand signal. F43 capped that recent spike to 2.0× the L13 median in-place (window-limited: only attenuates spikes inside the recent window). Result: 48,816/26w → 64,080/26w with a smoother profile. Annotated in alert text so reviewers see exactly which historical index was capped.",
    "F44":     "Acct 23011, Mstyle FF25895 — re-classification trigger. Once F43 attenuated the recent spike, the model would have read the now-flat recent window as 'lumpy/zero-heavy' and routed to the wrong recipe. F44 re-runs classification on F43-cleaned history so the model sees the dense order pattern that survives — typically routing to a baseline-style recipe instead of Lumpy.",
    "F45":     "Defensive per-week guardrail. After F43+F44 reshape, F45 caps any individual forecast week so no single position can run away. Catches edge cases where seasonal positioning still amplifies a surviving spike past 2× the post-cap baseline. Records exactly how many weeks were clamped in the alert text for transparency.",
    "F46":     "Acct 23011, Mstyle FF25895 — post-F44 forecast rebuild. F46 rebuilds 26 weeks from the pre-disruption baseline (L26 nz-mean × seasonal shape) distributed smoothly. AI lands at ~1,380/wk × 26 ≈ 35,880 — matching planner intent of a steady cadence — instead of choppy 683/wk with mostly-zero weeks left over from F43+F44+F45 alone.",
    "F47":     "Acct 23011, Mstyle FF12660 — Walmart Seasonal Baseline. F47 (OOS rebuild-ramp normalization, 2026-05-07) detects ≥3 consecutive ship=0 weeks with ord>0 (active stockout) and caps each within-gap order at 1.3× the pre-OOS baseline. The first post-gap week is also capped if it lands ≥1.5× the baseline. Capped indices are passed to F39/F41 dedupe rules as `protected_indices` to prevent double-zeroing. Brought FF12660 from +27% over manual to -13% — the rebuild-order spikes (8,640/9,720/6,840 vs pre-OOS pace ~1,800) are normalized to true demand intent before downstream rules see them.",
    "F48":     "Acct 1864, Mstyle BB13437 — Amazon Seasonal Baseline. F48 (post-OOS spike-and-cooldown anchor, 2026-05-07) detects when L13 baseline is inflated by a recent rebuild-order spike followed by L4 cooling. Trigger A: max ord in L13 ≥ 2.5× median (excl max), spike in W-12..W-5, AND L4 all-weeks avg < L13 nz-avg × 0.80. Trigger B (Amazon): healthy POS AND L4 ord < POS_blend × 0.85. Action: cap baseline at MAX(L4_avg, L26_avg) × 1.20 (Amazon: MAX(L4, POS_blend) × 1.20). Brought BB13437 from +54% over manual to +3%, FF15592 from +33% to +24%. Critical detail: uses L4 ALL-weeks avg (zeros count as buyer-side pause signal), not L4 nz-avg.",
    "VP-Q6":   "OOS fill-rate calculation bug-fix (2026-05-07). The `Shpd_Wk_L13W_cust_` field is a per-week average, but the code was comparing it against the L13 ORDER total — off by a factor of 13. False-positived OOS for nearly every active record (e.g. BB13437 reported 9% fill-rate when true fill-rate was 120% catch-up). The cascade then forced L13 nz-avg baseline + enabled F13 drawdown lift + F38 trend lift, inflating cap_base by 50-100% across thousands of records. Fix: compare per-week-avg vs per-week-avg, raised OOS threshold from 0.85 → 0.70 to require a clearly broken fill-rate.",
    "F5":      "Acct 23011, Mstyle BB24711 — Walmart, PT_Item_Status contains 'NEW' AND zero recent quarter shipments. F5 routes through the new-item sibling-SKU logic (F1 mstyle-family rate × cust-scale × 0.5 conservative) instead of dropping into the standard Inactive bucket. Launching items inherit a reasonable forecast from their cousins.",
    "F8":      "A seasonal pattern is matched via Product_Category OR Product_Subcategory rather than only Master_Category. F8 broadens the seasonal coverage so items that are tagged at sub-category level (but not master) still get their seasonal shape applied. More items in the catalog get correct seasonal lifts.",
    "F19":     "An Inactive-classified item with manual_total ≥ 5,000 AND consumer POS > 0 (Amazon path) OR last non-zero order within 26w (non-Amazon). F19's --conservative-inactive (on by default 2026-05-06) applies a 50% manual-shaped floor instead of forcing AI to zero. Safety net for accounts where planner manuals are very high-confidence.",
    "F20":     "Acct 23011, Mstyle FF28402 — Walmart, planner manually zeroed all 26 weeks. F20 forces the Reactivating recipe output to zero too, even though the recipe could have synthesized a forecast from history. When the planner says 'this is dead', the AI agrees rather than fighting the explicit zero.",
    "S6":      "Acct 13640, Mstyle FF15585 — Burlington Coat Factory, model 'Inactive+S6 (off-price)'. S6 detects small-chain off-price retailers (Burlington, Big Lots, Five Below, etc.) that buy in lumps with long quiet periods. AI lands at 120 units vs 3,400 manual (-96%) — the manual flat plan was pretending Burlington orders weekly, but they don't.",
    "M1":      "Acct 23011, Mstyle BB13435CLR/12 — M1 is the L52 ceiling: caps the forecast at the long-run weekly rate × 26 when other lifts (POS blend, F13 drawdown, F15 buyer premium) compound. Without M1, stacked lifts could drift the 26-week total to 2-3× the long-run pace; M1 anchors it to history.",
    "F1":      "Acct 23011, Mstyle BB0083 — Walmart, zero order history but model 'Inactive (zero order history)'. F1 (mstyle-family fallback) computes the median weekly rate across all sibling SKUs sharing this Mstyle at OTHER customers, scales it by a customer-specific factor and a 0.5 conservative multiplier, and uses that as the rate.",
}

# ─────────────────────────────────────────────────────────────────────────────
slides = []

# ── 1. Title ──────────────────────────────────────────────────────────────────
def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    add_rect(s, 0, Inches(2.5), SLIDE_W, Inches(0.05), fill=TEAL)
    add_rect(s, 0, Inches(5.4), SLIDE_W, Inches(0.05), fill=GOLD)
    add_text(s, MARGIN, Inches(2.7), CONTENT_W, Inches(1.1),
             "Inventory Forecaster", size=44, bold=True, color=WHITE)
    add_text(s, MARGIN, Inches(3.7), CONTENT_W, Inches(0.7),
             "How it works — in plain English (33 refinements)", size=22, color=GOLD)
    add_text(s, MARGIN, Inches(4.5), CONTENT_W, Inches(0.6),
             "A guide for inventory planners. No math required.",
             size=14, color=GREY_LT)
    add_text(s, MARGIN, Inches(6.6), CONTENT_W, Inches(0.4),
             "Pets+People Planning",
             size=10, color=GOLD)
    return s
slides.append(slide_title)

# ── 2. The big idea ──────────────────────────────────────────────────────────
def slide_big_idea():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "What this tool actually does",
                 "The big picture in three sentences")
    # Three big cards
    cards = [
        ("READS", NAVY,
         "It pulls one full year of order history for every active product/customer combo from Quickbase. "
         "For Amazon items, it also pulls Amazon's consumer sales data so we can see what shoppers are actually buying."),
        ("PREDICTS", TEAL,
         "For each combo, it picks the right kind of math (one of 9 'recipes') based on whether the demand is "
         "steady, lumpy, ramping up, or quiet. Then it predicts how many units the customer will order each week "
         "for the next 26 weeks (6 months)."),
        ("WRITES & FLAGS", ORANGE,
         "It writes the AI prediction back into Quickbase as 'AI_PRJ_W1' through 'AI_PRJ_W26', and adds a "
         "plain-English alert if the AI number differs from your manual projection by more than 5%. "
         "You stay in control — the AI just tells you where to look."),
    ]
    for i, (h, c, body) in enumerate(cards):
        top = Inches(1.20 + i * 1.95)
        add_rect(s, MARGIN, top, CONTENT_W, Inches(1.75), fill=WHITE, line=GREY_LT)
        add_rect(s, MARGIN, top, Inches(2.3), Inches(1.75), fill=c)
        add_text(s, MARGIN, top, Inches(2.3), Inches(1.75),
                 h, size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, MARGIN + Inches(2.5), top + Inches(0.15),
                 CONTENT_W - Inches(2.6), Inches(1.55),
                 body, size=13, color=GREY_DK, anchor=MSO_ANCHOR.MIDDLE)
    page_footer(s, 2, 35)
    return s
slides.append(slide_big_idea)

# ── 3. Pipeline at a glance ───────────────────────────────────────────────────
def slide_pipeline():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "How a forecast happens — step by step",
                 "From clicking 'run' to seeing the AI numbers in Quickbase")
    steps = [
        ("Step 1 — Gather data",
         "Pull all 52 weeks of past orders, plus consumer-sales data for Amazon items "
         "and the case-pack size for every SKU.",
         NAVY),
        ("Step 2 — Pick a recipe",
         "Look at the order pattern. Steady weekly orders? Lumpy and unpredictable? "
         "Brand new with little history? Each pattern gets its own forecasting recipe.",
         TEAL),
        ("Step 3 — Predict 26 weeks",
         "Run the chosen recipe to produce a number for each of the next 26 weeks. "
         "Bump up Prime Day weeks for Amazon, and the Fall Deal weeks for everyone.",
         ORANGE),
        ("Step 4 — Sanity check & save",
         "Round to case-pack multiples. If the customer already has a confirmed PO in a future week, "
         "set that AI week to zero (don't double-count). Save back to Quickbase + write an alert.",
         GOLD),
    ]
    box_h = Inches(1.30)
    box_w = Inches(11.50)
    top_y = Inches(1.20)
    for i, (head, body, c) in enumerate(steps):
        top = top_y + Inches(i * 1.45)
        add_rect(s, MARGIN + Inches(0.5), top, box_w, box_h, fill=WHITE, line=GREY_LT)
        add_rect(s, MARGIN + Inches(0.5), top, Inches(2.7), box_h, fill=c)
        add_text(s, MARGIN + Inches(0.5), top, Inches(2.7), box_h,
                 head, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, MARGIN + Inches(3.3), top + Inches(0.10),
                 box_w - Inches(2.9), box_h - Inches(0.20),
                 body, size=12, color=GREY_DK, anchor=MSO_ANCHOR.MIDDLE)
    page_footer(s, 3, 35)
    return s
slides.append(slide_pipeline)

# ── 4. The 9 recipes (model classification) ───────────────────────────────────
def slide_classify():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "The 9 forecasting recipes — and when each one is used",
                 "Plain-English version of how the tool picks a recipe")
    rows = [
        ("Customer hasn't ordered in 6+ months",
         "Inactive",
         "Forecast = 0. We don't invent demand from thin air."),
        ("Brand-new item ('NEW' status, no order history)",
         "Pre-launch",
         "Use planner's manual projection. Tool has nothing to compare to."),
        ("Amazon Private Label (Amazon-branded, we make)",
         "Normal Amazon recipe",
         "We ship these. Order/ship history + Amazon POS available — same recipe selection as any other Amazon item with POS blend."),
        ("Off-price retailer (Burlington, Ross, Big Lots, etc.)",
         "One-Time-Buy",
         "These accounts buy in occasional big lumps. Forecast zero between events."),
        ("Spotty history (lots of zero weeks)",
         "Sparse",
         "Average the weeks where orders did happen, scale by season."),
        ("Item just came back from a quiet period",
         "Reactivating",
         "Use the post-restart average; respect any recent ramp."),
        ("Steady weekly orders (most pet/CPG items)",
         "Steady-Demand",
         "Smooth out noise; use last 13 weeks as the strongest signal."),
        ("Up-and-down lumpy ordering",
         "Lumpy",
         "Predict order size separately from how often orders come."),
        ("Smooth dense orders with seasonal pattern",
         "Seasonal Baseline",
         "Anchor on recent average; adjust for time of year + Amazon shopper data."),
    ]
    add_table(s, MARGIN, Inches(1.05), CONTENT_W, Inches(5.6),
              ["What the order pattern looks like", "Recipe name", "What it does"],
              rows, size=11, hdr_size=12)
    add_text(s, MARGIN, Inches(6.75), CONTENT_W, Inches(0.30),
             "Routing is automatic — the tool reads the pattern and assigns a recipe. "
             "You don't pick; you just see which one was used in the alert.",
             size=10, color=GREY_MD)
    page_footer(s, 4, 35)
    return s
slides.append(slide_classify)

# ── 5. The four primary recipes — explained simply ────────────────────────────
def slide_models():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "The four main recipes — what they do, simply",
                 "These cover most of your active products")
    cards = [
        ("Steady-Demand", TEAL,
         "Items that ship a similar amount every week",
         ["Looks at the past year, but weights the most recent 3 months 3x heavier",
          "Smooths out random week-to-week noise",
          "Keeps a slight trend (climbing or sliding) if one is real",
          "Like predicting tomorrow's temperature from this week's pattern"]),
        ("Lumpy", ORANGE,
         "Orders are big, but with quiet weeks in between",
         ["Predicts two things separately: how big each order is + how often they come",
          "If orders come every 3 weeks at ~500 units, it spaces them out that way",
          "Won't smear demand into weeks the customer never orders",
          "Like predicting bus arrivals — not every minute, but the rhythm"]),
        ("Seasonal Baseline", NAVY,
         "Steady items where time-of-year matters",
         ["Anchored on the average of weeks where orders actually happened (last 3 mo)",
          "Blended with Amazon's consumer-sales rate for Amazon items",
          "Time-of-year shape is dampened so one big week doesn't blow up the forecast",
          "Prime Day and Fall Deal lifts are added on top"]),
        ("Sparse", GOLD,
         "Less than 3 months of order history in a year",
         ["Averages the weeks where the customer actually ordered something",
          "Applies the year-shape so seasonal items still spike at the right time",
          "Caps the total at a sensible ceiling so one big order doesn't dominate",
          "Like estimating restaurant traffic from 5 weeks of receipts"]),
    ]
    grid = [(0, 0), (0, 1), (1, 0), (1, 1)]
    cw = Inches(6.20); ch = Inches(2.85)
    for (row, col), (name, color, sub, bullets) in zip(grid, cards):
        x = MARGIN + Inches(col * 6.30)
        y = Inches(1.05 + row * 3.00)
        add_rect(s, x, y, cw, ch, fill=WHITE, line=GREY_LT)
        add_rect(s, x, y, cw, Inches(0.55), fill=color)
        add_text(s, x + Inches(0.20), y, cw - Inches(0.4), Inches(0.55),
                 name, size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.20), y + Inches(0.55),
                 cw - Inches(0.4), Inches(0.35),
                 sub, size=10, color=GREY_MD, anchor=MSO_ANCHOR.MIDDLE)
        add_bullets(s, x + Inches(0.20), y + Inches(0.95),
                    cw - Inches(0.4), Inches(1.85),
                    bullets, size=10, bullet_color=color)
    page_footer(s, 5, 35)
    return s
slides.append(slide_models)

# ── 6. Event calendar ─────────────────────────────────────────────────────────
def slide_events():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Two big seasonal events the tool handles",
                 "Why your AI forecast 'leans in' during certain weeks")
    # Two event "cards" — gives the long "why" text room to breathe instead
    # of squeezing it into a narrow table cell
    events = [
        ("Prime Day pre-buy", NAVY,
         "Weeks 5–9   ·   Amazon only   ·   Lifts orders 10–25%",
         "Schedule: W5=1.10, W6=1.15, W7=1.25, W8=1.25, W9=1.20.",
         "Amazon orders product 6–8 weeks BEFORE the July consumer event. The 'pre-buy' is what we "
         "ship — not the consumer sale itself."),
        ("Fall Deal pre-buy", ORANGE,
         "Weeks 23–25   ·   Every retailer   ·   Lifts orders 12%",
         "Flat lift × 1.12 across all three weeks.",
         "Smaller pre-order cycle for fall promotions. Not Amazon-specific — most retailers run "
         "end-of-summer / back-to-school deals."),
    ]
    for i, (name, c, line, sched, why) in enumerate(events):
        top = Inches(1.05 + i * 1.45)
        add_rect(s, MARGIN, top, CONTENT_W, Inches(1.30), fill=WHITE, line=GREY_LT)
        add_rect(s, MARGIN, top, Inches(2.6), Inches(1.30), fill=c)
        add_text(s, MARGIN, top, Inches(2.6), Inches(1.30),
                 name, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, MARGIN + Inches(2.75), top + Inches(0.08),
                 CONTENT_W - Inches(2.85), Inches(0.30),
                 line, size=12, bold=True, color=NAVY)
        add_text(s, MARGIN + Inches(2.75), top + Inches(0.40),
                 CONTENT_W - Inches(2.85), Inches(0.28),
                 sched, size=10, color=GREY_MD)
        add_text(s, MARGIN + Inches(2.75), top + Inches(0.72),
                 CONTENT_W - Inches(2.85), Inches(0.55),
                 why, size=11, color=GREY_DK)
    add_text(s, MARGIN, Inches(4.10), CONTENT_W, Inches(0.45),
             "How the tool knows it's an Amazon record", size=14, bold=True, color=NAVY)
    add_bullets(s, MARGIN, Inches(4.50), CONTENT_W, Inches(2.5), [
        "It checks if the customer name contains the word 'AMAZON'",
        "If yes: pull Amazon shopper data and apply Prime Day lift",
        "If no: skip both — most other retailers don't share consumer sales with us anyway",
        "Online-only retailers (Chewy, Petco.com, PetSmart.com) get a special 'recent acceleration' check instead",
    ], size=12)
    page_footer(s, 6, 35)
    return s
slides.append(slide_events)

# ─────────────────────────────────────────────────────────────────────────────
# Sales Index / Seasonality section (slides 7-11)
# ─────────────────────────────────────────────────────────────────────────────
import json as _json
import os as _os
_DERIVED_PATH = _os.path.join(_os.path.dirname(_os.path.abspath(__file__)),
                              "scripts", "derived_category_profiles.json")

def _load_subcat_data():
    """Pull qualifying subcategory profiles for the deck table."""
    if not _os.path.exists(_DERIVED_PATH):
        return []
    with open(_DERIVED_PATH) as f:
        d = _json.load(f)
    rows = []
    for key, payload in (d.get("by_subcategory") or {}).items():
        cat, _, subcat = key.partition("||")
        stats = payload.get("stats") or {}
        n = stats.get("consistent_skus", 0) or 0
        prof = payload.get("profile") or []
        if not prof or len(prof) != 12:
            continue
        # Only show qualifying ones (SKU gate >10) per the live forecaster rule
        if n <= 10:
            continue
        rows.append({"cat": cat, "sub": subcat, "n": n, "prof": prof})
    rows.sort(key=lambda r: (r["cat"], r["sub"]))
    return rows

_SUBCAT_ROWS = _load_subcat_data()

# ── 7. Sales Index — what it is, how it's matched ─────────────────────────────
def slide_index_intro():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Sales Index — how time-of-year shaping works",
                 "Forecasts get bent by season; here's how the right shape gets picked")
    add_text(s, MARGIN, Inches(1.00), CONTENT_W, Inches(0.45),
             "What is a 'sales index'?", size=14, bold=True, color=NAVY)
    add_text(s, MARGIN, Inches(1.40), CONTENT_W, Inches(0.95),
             "A sales index is a 12-number profile (one per month) that says how much busier or quieter "
             "a product is each month relative to its yearly average. Charcoal in June might be 2.05 "
             "(twice as busy as average); charcoal in December might be 0.20 (one-fifth of average). "
             "We multiply the AI baseline by these numbers to bend the forecast in or out by season.",
             size=12, color=GREY_DK)
    add_text(s, MARGIN, Inches(2.55), CONTENT_W, Inches(0.45),
             "Match priority — the tool tries these in order, takes the first hit:",
             size=14, bold=True, color=NAVY)
    chain = [
        ("1", NAVY, "Planner-curated Season tag",
         "If the SKU has a Season tag in Quickbase Styles (Holiday, Halloween, July 4th, etc.) we use "
         "that hand-tuned profile first. Highest trust — the planner explicitly picked it."),
        ("2", TEAL, "Data-derived (Category + Subcategory)",
         "We've analyzed 2024-2026 invoice ship history and built a profile for every Category + "
         "Subcategory pair that has more than 10 SKUs (statistically reliable). Most specific shape match wins."),
        ("3", ORANGE, "Data-derived (Category alone)",
         "If the subcategory match doesn't qualify (<11 SKUs), fall back to the broader category profile "
         "if it qualifies. Same rules — built from the same 3-year invoice history."),
        ("4", GOLD, "Hand-curated keyword fallback",
         "For items missing structured tags, we keyword-match against the description: 'charcoal', "
         "'mosquito', 'sunscreen', 'ice melt', etc. Hand-curated profiles for known seasonal categories."),
    ]
    top = Inches(3.10)
    for n, c, head, body in chain:
        add_rect(s, MARGIN, top, CONTENT_W, Inches(0.92), fill=WHITE, line=GREY_LT)
        add_rect(s, MARGIN, top, Inches(0.55), Inches(0.92), fill=c)
        add_text(s, MARGIN, top, Inches(0.55), Inches(0.92),
                 n, size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, MARGIN + Inches(0.70), top + Inches(0.05),
                 CONTENT_W - Inches(0.75), Inches(0.30),
                 head, size=12, bold=True, color=NAVY)
        add_text(s, MARGIN + Inches(0.70), top + Inches(0.35),
                 CONTENT_W - Inches(0.75), Inches(0.55),
                 body, size=11, color=GREY_DK)
        top += Inches(0.97)
    # Data-hygiene callout — explicit answer to "did you parse out outliers / OTBs / OOS?"
    add_rect(s, MARGIN, Inches(6.85), CONTENT_W, Inches(0.55), fill=WHITE, line=GREY_LT)
    add_text(s, MARGIN + Inches(0.10), Inches(6.88), CONTENT_W - Inches(0.20), Inches(0.22),
             "Data hygiene before profiles are built (handled in build_category_profiles_from_report.py):",
             size=9, bold=True, color=NAVY)
    add_text(s, MARGIN + Inches(0.10), Inches(7.10), CONTENT_W - Inches(0.20), Inches(0.30),
             "OTBs / promo blasts excluded (SKU must ship ≥10 months over ≥12-month span, ≥50% active rate)  •  "
             "OOS months dropped (May-Sep 2025 tariff window)  •  outliers clamped (0.10x ↔ 4.00x, mean=1.0)  •  "
             "≥3 consistent SKUs and ≥100k units per category  •  2024 weighted 2× (cleanest year)",
             size=8, color=GREY_DK)
    page_footer(s, 7, 35)
    return s
slides.append(slide_index_intro)

# ── 8. Planner-curated Season tags table ──────────────────────────────────────
SEASON_PROFILES_DECK = [
    ("Holiday",          "Thanksgiving / Christmas paper goods", "Aug-Nov", "Oct-Nov peak (2.30, 2.20)"),
    ("Halloween",        "Halloween category", "Jul-Sep",  "Aug peak (2.40)"),
    ("July 4th",         "Independence Day",   "Apr-Jun",  "May peak (2.20)"),
    ("Easter",           "Easter category",    "Jan-Mar",  "Mar peak (2.10)"),
    ("Valentines Day",   "Valentine's gifting","Nov-Jan",  "Dec peak (2.15)"),
    ("St Patrick's Day", "St Pat's category",  "Dec-Feb",  "Jan peak (2.00)"),
    ("Pride",            "June Pride category","Mar-May",  "Apr peak (2.00)"),
    ("Spring/Summer",    "Outdoor lifestyle",  "Feb-Jun",  "Mar-May peak (1.65, 1.55)"),
    ("Fall/Winter",      "Cold-weather indoor","Aug-Dec",  "Sep-Oct peak (1.65, 1.65)"),
]
def slide_season_tags():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Priority 1 — Planner-curated Season tags",
                 "9 hand-tuned profiles applied when the SKU has a 'Season' tag in QB Styles")
    add_text(s, MARGIN, Inches(1.00), CONTENT_W, Inches(0.45),
             "How the orderingweeks were derived", size=13, bold=True, color=NAVY)
    add_text(s, MARGIN, Inches(1.40), CONTENT_W, Inches(0.55),
             "Retail ordering leads consumer demand by 4-8 weeks. Peaks in these tables reflect when "
             "RETAILERS PLACE ORDERS, not when consumers buy. Easter peaks in Feb-Mar (not April).",
             size=11, color=GREY_DK)
    add_table(s, MARGIN, Inches(2.10), CONTENT_W, Inches(4.5),
              ["Tag", "Use case", "Order window (months)", "Order peak"],
              SEASON_PROFILES_DECK, size=11, hdr_size=12)
    add_text(s, MARGIN, Inches(6.75), CONTENT_W, Inches(0.30),
             "Tags applied automatically — no SKU gate (planner-trusted). Floor still at 1.0 (no demand cuts).",
             size=10, color=GREY_MD)
    page_footer(s, 8, 35)
    return s
slides.append(slide_season_tags)

# ── 9-10. Data-derived subcategory profile tables (split across 2 slides) ────
MONTHS_SHORT = ["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]

def _format_subcat_row(r):
    p = r["prof"]
    peak_idx = p.index(max(p))
    # Show only month indexes where the floored profile lifts demand (>1.0)
    lift_months = [MONTHS_SHORT[i] for i, v in enumerate(p) if v > 1.05]
    lift_str = ", ".join(lift_months) if lift_months else "—"
    return [r["cat"], r["sub"], str(r["n"]),
            f"{MONTHS_SHORT[peak_idx]} ({p[peak_idx]:.2f}x)",
            lift_str]

def slide_subcat_table(part_idx, rows, page_n):
    s = prs.slides.add_slide(BLANK)
    slide_header(s, f"Priority 2 — Data-derived subcategory profiles ({part_idx} of 2)",
                 "Built from 2024-2026 invoice ship history — 35 qualifying subcategories")
    table_rows = [_format_subcat_row(r) for r in rows]
    add_table(s, MARGIN, Inches(1.00), CONTENT_W, Inches(5.95),
              ["Category", "Subcategory", "SKUs", "Peak month (index)", "Months that lift demand (>1.05x)"],
              table_rows, size=9, hdr_size=10)
    add_text(s, MARGIN, Inches(7.05), CONTENT_W, Inches(0.20),
             "SKUs = number of consistent SKUs in that subcategory. Floor=1.0 — only LIFTS, never cuts.",
             size=9, color=GREY_MD)
    page_footer(s, page_n, 35)
    return s

# Split qualifying rows across 2 slides
_HALF = (len(_SUBCAT_ROWS) + 1) // 2
slides.append(lambda rows=_SUBCAT_ROWS[:_HALF]:
              slide_subcat_table(1, rows, 9))
slides.append(lambda rows=_SUBCAT_ROWS[_HALF:]:
              slide_subcat_table(2, rows, 10))

# ── 11. Hand-curated keyword fallback ─────────────────────────────────────────
KEYWORD_GROUPS = [
    ("Outdoor cooking / grilling",   "charcoal · chimney · fire starter · lighter fluid · grill brush · grill cleaner · Kingsford · wooden fire", "Apr-Aug",  "Jun (2.05x)"),
    ("Pest control / repellent",     "mosquito · insect repel · bug repel",                                                                      "May-Sep",  "Jul (2.05x)"),
    ("Sun care",                     "sunscreen · sun care · sunblock",                                                                          "May-Aug",  "Jun-Jul (2.05x)"),
    ("Home fragrance / air-care",    "air freshener · deodorizing ball · scent booster · Fraganzia",                                             "Apr-Aug",  "Jun (1.35x)"),
    ("Outdoor party disposables",    "snack bowl · paper bowl · paper plate · paper cup",                                                        "Apr-Aug",  "Jun (1.65x)"),
    ("Holiday / Christmas",          "holiday · christmas",                                                                                      "Sep-Jan",  "Dec (2.50x)"),
    ("Ice melt / de-icer",           "ice melt · de-icer",                                                                                       "Nov-Feb",  "Dec (2.00x)"),
]
def slide_keyword_fallback():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Priority 4 — Keyword fallback profiles",
                 "Hand-curated patterns for items that don't have structured category tags")
    add_text(s, MARGIN, Inches(1.00), CONTENT_W, Inches(0.55),
             "When an item has no Season tag and its category/subcategory hasn't qualified for a "
             "data-derived profile, the tool checks the description text against these keyword groups. "
             "First match wins. Each group has its own 12-month profile floored at 1.0.",
             size=11, color=GREY_DK)
    add_table(s, MARGIN, Inches(1.85), CONTENT_W, Inches(4.5),
              ["Group", "Keywords (case-insensitive)", "Order window", "Order peak"],
              KEYWORD_GROUPS, size=9, hdr_size=11)
    add_text(s, MARGIN, Inches(6.55), CONTENT_W, Inches(0.45),
             "Why a fallback?", size=13, bold=True, color=NAVY)
    add_text(s, MARGIN, Inches(6.95), CONTENT_W, Inches(0.30),
             "Some items have legacy / sparse category data. Keyword match catches them by description text "
             "(e.g. 'Kingsford 16 lb Charcoal Bag' → grill profile).",
             size=10, color=GREY_MD)
    page_footer(s, 11, 35)
    return s
slides.append(slide_keyword_fallback)

# ── 12. The 33 refinements — friendly intro ───────────────────────────────────
def slide_ref_intro():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "We tuned the tool 29 times. Here's why.",
                 "Real-world buying behavior is messier than any one formula handles")
    add_text(s, MARGIN, Inches(1.05), CONTENT_W, Inches(0.45),
             "When the tool first launched, we noticed things like:", size=14, color=GREY_DK)
    add_bullets(s, MARGIN + Inches(0.30), Inches(1.50), CONTENT_W - Inches(0.30),
                Inches(2.2), [
        "Forecasts dropped off after Prime Day because the post-event 'quiet weeks' looked like demand died",
        "Off-price retailers (Burlington, Big Lots) got phantom weekly forecasts when they actually buy in lumps",
        "Items in real decline kept getting full-strength forecasts — no decline detection",
        "Amazon's confirmed POs were already in the system, but the AI added MORE on top → over-projection",
        "Heavy-seasonal items (charcoal, grill brushes) got their seasonal shape softened too much",
    ], size=12)
    add_text(s, MARGIN, Inches(3.85), CONTENT_W, Inches(0.45),
             "The 33 refinements address those patterns. They fall into three families:",
             size=14, color=GREY_DK)
    cards = [
        ("Big-picture rules", NAVY, "4 rules",
         "From the VP of Planning. Highest impact: how to define 'normal demand', "
         "how to handle stockouts, when to enforce ordering rhythms, and how to handle confirmed POs."),
        ("Customer-type overrides", ORANGE, "8 rules",
         "Different retailers buy differently. Off-price chains, international wholesalers, "
         "online-only stores, lumpy big-volume buyers — each gets its own treatment."),
        ("Math fine-tuning", TEAL, "17 rules",
         "Smaller fixes for specific patterns: thin order history, recent slowdowns, "
         "stuck cadences, drawdown detection, end-of-life items."),
    ]
    cw = Inches(4.10); ch = Inches(2.30)
    for i, (title, c, count, body) in enumerate(cards):
        x = MARGIN + Inches(i * 4.20)
        y = Inches(4.55)
        add_rect(s, x, y, cw, ch, fill=WHITE, line=GREY_LT)
        add_rect(s, x, y, cw, Inches(0.55), fill=c)
        add_text(s, x, y, cw, Inches(0.55), title, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, y + Inches(0.65), cw, Inches(0.35), count,
                 size=12, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.20), y + Inches(1.05),
                 cw - Inches(0.4), ch - Inches(1.2),
                 body, size=11, color=GREY_DK, anchor=MSO_ANCHOR.TOP)
    page_footer(s, 12, 35)
    return s
slides.append(slide_ref_intro)

# ── 8–11. The 4 big-picture rules (VP-Q stack) — plain language ──────────────
VP_Q_SLIDES = [
    ("VP-Q1", "How we define 'normal demand'",
     "We averaged ALL 13 weeks of recent history — including the quiet weeks right after Prime Day. "
     "Those quiet weeks aren't a real demand drop; they're just leftover Prime Day inventory burning off. "
     "Including them dragged our 'normal' baseline down 15-20%.",
     "Now we average ONLY the weeks where orders actually happened — the quiet 'drawdown' weeks are skipped. "
     "If we don't have enough non-zero weeks to use, we fall back to looking at 6 months instead of 3.",
     "Your AI forecast better reflects steady-state demand — not a number suppressed by a normal post-event lull. "
     "On account 1864, this alone added 18.5% to total predicted demand because it stopped under-projecting.",
     "The alert text always shows you which 'baseline mode' was used."),
    ("VP-Q2", "Treating stockouts like the demand they really were",
     "When Amazon's order got cancelled because we were out of stock, the system saw '0 orders' for that week — "
     "as if customers didn't want it. That artificially shrank our forecast for the next time.",
     "We now pull the cancellation reason from Order History. If it was OOS-driven (inventory error, supplier "
     "delay), we keep the original order quantity as 'demand intent'. If it was demand-driven (customer "
     "cancellation, future delete), we treat it as zero. Ambiguous cases stay as-is.",
     "Items that have a history of stockouts will now get correctly-sized forecasts going forward — we won't "
     "punish a SKU because OUR stock issue made it look like demand died.",
     "Off by default; gets turned on automatically during back-tests."),
    ("VP-Q3", "Stop forcing fake bi-weekly patterns",
     "If the system saw a few quiet weeks in the recent past, it would assume the customer orders every other "
     "week and start INSERTING zeros into the forecast — even when the latest pattern was weekly again.",
     "We only enforce a rhythm if the gaps are at least 3 weeks long AND consistent (monthly, every 5 weeks, etc.). "
     "Anything shorter just gets smoothed across weeks without forced zeros.",
     "Forecasts no longer alternate between 'big' and 'zero' weeks for items that quietly returned to weekly "
     "ordering. You'll see fewer 'why is the AI predicting zero?' moments.",
     "Lumpy-recipe items that hit this rule get smoothed to weekly automatically."),
    ("VP-Q4", "Don't double-count confirmed customer POs",
     "If a customer already had a confirmed open PO for week 3, the AI was projecting MORE units on top — "
     "as if the PO didn't exist. The replenishment system was then double-counting those units.",
     "For each item, we pull every open PO with a customer ship-by date. Any AI forecast week that falls "
     "in the same week as a confirmed PO is set to ZERO — the PO IS the demand for that week.",
     "AI forecasts now match what your replenishment system thinks. Stops the chronic over-ordering in the "
     "first 4–6 forward weeks where most confirmed POs land.",
     "On by default. Turn off with --no-po-zero only if you're back-testing the old behavior."),
]
# Consolidated VP-Q rendering — 2 cards per page, same card style as F-rules
# (refinement_card with examples). Footnote folded into the WHY section so
# the card stays self-contained.
def slide_vp_group_page(page_in_group, total_pages, page_items, abs_page_n):
    s = prs.slides.add_slide(BLANK)
    sub = (f"VP-led changes (page {page_in_group} of {total_pages})"
           if total_pages > 1 else "VP-led changes — plain English")
    slide_header(s, "Big-picture rules — the VP-Q stack", sub)
    top = Inches(0.85)
    card_h = Inches(2.85)  # 2 cards × 2.85 + 0.05 gutter + 0.85 top = 6.60"
    for (code, title, before, after, why, foot) in page_items:
        # Footnote dropped — example below carries the same intent more concretely.
        refinement_card(s, MARGIN, top, CONTENT_W, card_h, code, title,
                        before, after, why, example=EXAMPLES.get(code, ""))
        top += card_h + Inches(0.05)
    page_footer(s, abs_page_n, 35)
    return s

# Chunk into pages of 2 cards each
_vp_pages = [VP_Q_SLIDES[i:i+2] for i in range(0, len(VP_Q_SLIDES), 2)]
for _pi, _chunk in enumerate(_vp_pages):
    slides.append(lambda chunk=_chunk, pi=_pi+1, tot=len(_vp_pages), p=8+_pi:
                  slide_vp_group_page(pi, tot, chunk, p))

# ── 12-14. R-series — customer-type overrides — plain language ───────────────
R_REFINEMENTS = [
    ("R1", "Off-price retailers buy in lumps, not weekly",
     "Burlington, Ross, TJ Maxx, Big Lots, Five Below, Kohl's, etc. were getting weekly forecasts based "
     "on a sparse history with 1-3 big orders.",
     "Detect the lump pattern (a few big orders separated by long gaps). Set the recipe to 'One-Time-Buy' — "
     "forecast zero between events.",
     "No more phantom W1 demand for accounts that buy 3 times a year and would never order weekly."),
    ("R2", "Cap on lumpy retailers — sanity ceiling",
     "Sparse-recipe items at off-price retailers were over-predicting by 38% — one big historical order "
     "amplified by seasonal scaling.",
     "Cap the total 26-week forecast at 1.5x what they'd buy at their average pace over 6 months.",
     "Forecasts for these accounts can't run away even if one outlier order hits the math."),
    ("R3", "Inactive items with real residual demand",
     "Items routed 'Inactive' had forecast = 0, but some had small but real activity over the prior 6-12 months.",
     "If the customer ordered at least 4 times in the last 6 months and 8 times in the last year, "
     "apply a small floor instead of zero.",
     "Catches the 'slow steady seller' that wouldn't otherwise show up in the AI."),
    ("R5", "International retailers order seasonally, not weekly",
     "Petbarn (Australia), Loblaws (Canada), Mexican retailers ('Comercializadora', 'Grup') were marked "
     "'Inactive' because they hadn't ordered in 13 weeks.",
     "For these accounts, only mark Inactive if there's been NO order in 26 weeks (not 13).",
     "International retailers come back; we shouldn't write them off after one quiet quarter."),
    ("R6", "Lumpy recipe needed help on big steady items",
     "The Lumpy recipe was systematically under-predicting big stable items by ~511K units across 368 records.",
     "When recent order sizes are consistent and big enough (50+ per week), boost the predicted order size "
     "toward the actual recent average.",
     "Closes a known gap on high-volume retailers that order big regular orders."),
    ("R7", "Make sure Amazon items participate in Fall Deal",
     "If an Amazon item had no historical signal for the Fall Deal weeks, the forecast missed it entirely — "
     "no fall pre-buy at all.",
     "If we don't see anything in the Fall Deal weeks, we INSERT an extra order at +12% lift.",
     "Amazon items don't get skipped on the second-biggest event of the year."),
    ("R8", "When 2 big orders dominate a 'busy' history",
     "Items with 9 active weeks BUT where 2 huge orders dwarfed the others were getting an inflated 'average'.",
     "If the top 2 orders are more than 60% of the recent total, use the MEDIAN order size instead of the average.",
     "The median tells the truth about repeat-order size; the average gets hijacked by 2 burst weeks."),
    ("R9", "Hard ceiling on Reactivating-recipe items",
     "The Reactivating recipe was over-projecting by 24% on items coming back from a quiet period.",
     "Cap baseline at 2x the long-run weekly average over the past year.",
     "Recovering items get predicted, but never more than twice their historical pace."),
]
def slide_r_block(idx, items, page_n):
    s = prs.slides.add_slide(BLANK)
    code_list = ", ".join(c for c, *_ in items)
    slide_header(s, f"Customer-type overrides ({idx} of 3)",
                 f"R-series - {code_list}")
    # 3 cards x 2.10" + 2 x 0.05" gutter = 6.40"; top 0.85" = 7.25" content height (fits 7.30")
    card_h = Inches(2.10); card_w = CONTENT_W
    for i, (code, title, before, after, why) in enumerate(items):
        top = Inches(0.85 + i * 2.15)
        refinement_card(s, MARGIN, top, card_w, card_h, code, title,
                        before, after, why, example=EXAMPLES.get(code, ""))
    page_footer(s, page_n, 35)
    return s
for i in range(3):
    chunk = R_REFINEMENTS[i*3:(i+1)*3]
    slides.append(lambda chunk=chunk, idx=i+1, p=17+i:
                  slide_r_block(idx, chunk, p))

# ── 15-22. F-series — math fine-tuning, plain language ───────────────────────
F_GROUPS = [
    ("Math fine-tuning — recent vs. older history",
     [
        ("F4", "When the last 3 months are too thin to trust",
         "An item with only 3-4 active weeks in the last 13 was producing baselines near zero, "
         "even when the past 12 months looked healthy.",
         "If the recent window has fewer than 5 active weeks AND the past year has 8+, pull in "
         "the year-long average as a floor.",
         "Quiet quarters can't make a real product look dead."),
        ("F6 / F26 / F27", "Catching a recent slowdown — or pickup",
         "If the last 4 weeks were running half the rate of the last 13, the tool ignored it.",
         "If recent rate ≤ 50% of recent-quarter rate → cut baseline to 65%. "
         "If 50–70% → cut to 85%. "
         "If 30–60% ABOVE recent → bump baseline to 110% (a gentle ramp signal).",
         "The forecast leans into what you're seeing in the last month — not just the last quarter."),
     ]),
    ("Math fine-tuning — outliers & ceilings",
     [
        ("F25", "Drop one-off freak orders entirely",
         "A lone 744-unit order vs. a normal 50/week was dragging the average up wildly.",
         "If a single value is 5x the median AND there are 4+ other normal weeks, drop it from "
         "the average instead of trying to cap it.",
         "One weird order doesn't define the next 6 months."),
        ("F24", "Sanity ceiling on the final number",
         "After all the boosts (Amazon POS blend, drawdown lift, order-coverage lift) stacked up, "
         "the baseline could end up at 3-4x the recent rate.",
         "Cap the final baseline at L13 weekly avg × 2.0 (eased from 1.5× on 2026-05-06 for "
         "seasonal items). Profile and event lifts still apply on top.",
         "Hard upper limit so the math can't run away."),
     ]),
    ("Math fine-tuning — seasonality & events",
     [
        ("F7", "Use the seasonal peak, not the trough",
         "A grill brush in week 5 was anchored on its winter trough demand (low) — not its summer peak (high) "
         "where it was actually heading.",
         "If the category has a known seasonal pattern AND the year shows a strong peak, anchor on the "
         "historical peak adjusted for current week's seasonal factor.",
         "Items entering peak season get realistic forecasts, not ones based on the dead winter rate."),
        ("F11", "Prime Day lift shape — ramp up, then taper",
         "We were applying a flat +25% lift across weeks 7–9. Real pre-buys ramp up and taper off.",
         "Weeks 5/6/7/8/9 = +10%, +15%, +25%, +25%, +20%. Matches what planners actually see.",
         "AI forecast looks like a real Prime Day pattern, not a step function."),
        ("F16 / F16b", "Don't soften strong seasonal items too much",
         "Heavy-seasonal items (charcoal, fire starters, grill brushes) had their pattern dampened to ±20% — "
         "which crushed their natural 3x summer peak.",
         "If the category is known-seasonal AND the peak/trough ratio is steep enough AND volume is real "
         "(50+/week), let the raw seasonal shape through.",
         "Strong-seasonal items get their natural shape; everyone else stays softened."),
     ]),
    ("Math fine-tuning — Amazon & online retailers",
     [
        ("F13", "Drawdown means a refill is coming (Amazon)",
         "Amazon was selling product to consumers FASTER than they were ordering from us — they were burning "
         "down inventory. The order-side lull made AI under-forecast the upcoming refill.",
         "If consumer sales rate is meaningfully higher than the order rate AND consumer sales are healthy, "
         "raise the baseline to account for the inevitable replenishment.",
         "Catches Amazon's 'drawdown then big refill' pattern."),
        ("F14a / F14b", "Don't kill an item that's just having a buyer-side pause",
         "F10 (decline detection) was scaling Amazon items DOWN whenever buyer-side ordering paused — even "
         "when consumers were still buying through fine.",
         "If consumer sales stay strong (recent month is 50%+ of recent quarter, AND volume is meaningful), "
         "skip the decline-detection scale-down. Trust consumer sales over a buyer-side hiccup.",
         "Items don't get prematurely killed by a 1-2 week ordering lull."),
        ("F15", "Some accounts buy a bit more than they sell",
         "Some Amazon items showed planner orders consistently 30%+ above consumer sales rate — that's the "
         "buyer's choice for safety stock or distribution coverage.",
         "When the order/consumer-sales ratio is consistently >1.15 AND consumer sales are healthy, "
         "blend toward the higher order rate.",
         "AI respects the buyer's structural premium, not just consumer demand."),
        ("F18", "When Amazon shoppers are buying faster than orders imply",
         "Lumpy-recipe Amazon items had quantities that were too low because Amazon was selling through faster.",
         "If consumer sales are healthy AND running faster than the order pattern would predict, "
         "raise the predicted order size.",
         "Forecast catches up to actual consumer velocity."),
        ("T4", "Online-only retailers — Chewy, Petco.com, PetSmart.com",
         "These ecommerce stores have no consumer-sales feed for us, but planners see late-cycle ramps.",
         "If the last 4 weeks are running 5%+ above the last 13 weeks, blend baseline toward the recent "
         "rate (more weight if it's running 15%+ hot).",
         "Captures online retailers' fast acceleration without needing a POS feed."),
     ]),
    ("Math fine-tuning — declining & end-of-life items",
     [
        ("F10", "Detecting real declines (year-over-year check)",
         "Items in genuine decline were getting full-strength forecasts.",
         "Two checks must BOTH pass: recent month ≤ 70% of recent quarter AND year-over-year is also down. "
         "If both fire, scale the forecast down to the recent rate.",
         "The YoY check prevents firing on seasonal troughs that LOOK like decline but recover."),
        ("F22a / F22c", "Trailing zeros mean a drawdown — discount accordingly",
         "Items ending the recent history with 4-5 zero weeks in a row were mid-drawdown. The 'per-order' "
         "average was right, but the volume going forward should be lower.",
         "Discount baseline by trailing-zero count out of 13. Cap final baseline at recent average when "
         "the recent window is thin (less than 7 active weeks).",
         "Drawdown items get realistic next-quarter forecasts, not full-pace ones."),
        ("F23a / F23b", "Reactivating items — same drawdown logic + tame the seasonal shape",
         "Reactivating-recipe items on sparse history sometimes had a single big order multiplied 10x by "
         "seasonal positioning.",
         "Soften the seasonal shape (same as Seasonal Baseline). Apply trailing-zero discount up to 70%.",
         "No more wild seasonal multipliers from a single position-lucky order."),
     ]),
    ("Math fine-tuning — volume floors & sparse items",
     [
        ("F9", "Big-volume sparse items — use the MAX, not average",
         "Sparse items with 15K+ annual units sometimes had a quiet recent quarter that pulled the average down.",
         "Use whichever is highest: the recent 3-month, 6-month, or 12-month non-zero average.",
         "Strong-but-lumpy items don't get under-forecast by one quiet quarter."),
        ("F17 / F17b", "Get week-1 right for stuck cadences",
         "The Sparse recipe placed the next predicted order several weeks out — but the planner was expecting "
         "an order in week 1.",
         "Shift the cadence so week 1 gets seeded. Only do this if the item has enough volume (25+/week recent average).",
         "AI's first week aligns with what the planner sees coming. Low-volume tail items don't get inflated."),
        ("F28", "Lumpy recipe — don't go below recent average",
         "The Lumpy recipe was systematically running 15% below the recent quarter average across 414 records.",
         "Floor: don't let the Lumpy recipe output go below a fixed share of recent quarter average.",
         "Closes a known under-forecast bias on the Lumpy recipe."),
        ("F29", "New items with one shipping week",
         "A brand-new item that just started shipping might have only 1 active week in the last month — "
         "the original rule needed 2.",
         "Loosened to use any shipped week (last 4 or last 8) as the floor for the new-item forecast.",
         "Catches first-quarter shipping reality without over-restricting."),
     ]),
    ("Math fine-tuning — front-week guard rails",
     [
        ("F30", "Cap on big-volume Seasonal Baseline items",
         "Top-volume Seasonal Baseline records were running 13% over recent rate on average.",
         "Tighter cap on baseline-vs-recent-rate ratio for big items. Buyer-driven items have human plans "
         "we shouldn't out-forecast.",
         "Top items don't run hot from stacked lifts."),
        ("F30 (rev)", "Zero order history → zero forecast (always)",
         "Some recipes could synthesize forecasts from consumer sales alone on items with NO orders in 6+ months.",
         "Hard rule: if there are zero orders in the last 6 months, the AI forecast is zero — no exceptions.",
         "Customer hasn't ordered for half a year — POS alone isn't enough to project orders."),
        ("F31", "Two W1 problems in one rule",
         "(a) Brand new items with NEW status had bogus AI forecasts. "
         "(b) Steady-Demand recipe occasionally produced wild W1 numbers from trend extrapolation.",
         "(a) If status contains 'NEW' and zero history, pass through the planner's manual. "
         "(b) Cap W1 outliers when the median W1 across all items is +177% (sign that something's wrong).",
         "Pre-launch items don't get random predictions; W1 outliers don't blow up the front of the forecast."),
        ("F32", "Sparse items — per-week clamp",
         "The original cap was on the SUM of 26 weeks; rarely fired because sparse totals are small.",
         "Replace with a per-week clamp + a tiny-signal floor for items with near-zero recent history.",
         "Tighter, more reliable cap on the small sparse items."),
        ("F34", "New launches — don't penalize the ramp",
         "Items that started shipping ~26 weeks ago show 25+ weeks of zeros at the back of history. "
         "The decline-detection rule (F10) and the L52 ceiling (M1) read those pre-launch zeros as 'this "
         "item is dying' and squashed the forecast to near zero.",
         "Detect new launches: if weeks 27-51 ago add up to less than 1% of the last 26 weeks, mark as "
         "new launch. Skip F10 decline detection and the L52 ceiling so ramp-up volume is preserved.",
         "Newly-launched items get a real forecast based on the activity they DO have, not a zero based "
         "on the history they don't have yet."),
        ("F35", "Stockout backlog — strip pent-up orders",
         "When we're out of stock for several weeks the customer keeps re-ordering and qty piles up. "
         "On shipment resume the catch-up week showed 'base + backlog' and the tool over-projected.",
         "Detect 2-8 week zero gaps in dense ordering patterns. Use the planner's decay schedule "
         "(week 1 of stockout: 25% lost / 75% recoverable; week 2: 50%/50%; week 3: 75%/25%; week 4+: "
         "100% lost / 0% recoverable) to compute how much of the catch-up burst is legitimate backlog. "
         "Strip that amount from the post-gap catch-up weeks; the rest is left intact.",
         "Stockouts no longer poison the baseline. The forecaster sees true demand intent for the next "
         "26 weeks, not the artifact of a missed-shipment recovery."),
        ("F36", "Stock-up burn-off — wait until stores work through it",
         "After a big catch-up shipment, the customer is sitting on weeks of cover and won't re-order "
         "until POS sells through. Tool was projecting against post-shipment quiet weeks as decline.",
         "Amazon-only (POS data). Detect shipment cluster >=3x weekly POS rate, followed by quiet "
         "orders <0.5x POS rate. Weeks-of-supply = cluster qty / POS rate, minus weeks elapsed. "
         "Force AI W1..W{remaining-WOS} to 0 so cadence resumes at the right time.",
         "We stop projecting orders the customer can't possibly place. Once the stock-up burns through, "
         "the forecast picks back up exactly when the customer is ready to re-order."),
        ("F37", "Forward inventory shortfall — don't promise what we can't ship",
         "Forecast was projecting POs without checking next-26-week on-hand. If anticipated OH would "
         "run out, we projected shipments that physically can't happen and lost demand wasn't counted.",
         "Read Inv_Wk1..Inv_Wk26 (already AI-deducted). For each week: if OH < ship qty, cap to OH "
         "and roll unmet forward as a cohort that decays 25%/wk (lost at 4w+). Recovered backlog "
         "lands when a future week has spare OH.",
         "The AI forecast is now physically realistic given our inventory plan. Weeks where we'd "
         "stock out get capped, and the demand-loss math reflects the real-world consequence: the "
         "longer we miss, the more sales are gone for good."),
        ("F38", "Amazon POS-trend sensitivity — react to what shoppers are buying",
         "Amazon POS could swing 15-30% week-over-week with no change to forecast. ASINs that went "
         "'Not Buyable' or 'Suppressed' kept getting normal projections that wouldn't ship.",
         "Amazon-only. Compare L4w vs L13w POS. >+10% trend = lift baseline (F38b); but ignore if "
         "buybox dropped ≥10% below MAP (F38a, temp discount). <-10% = cut baseline (F38e); but "
         "ignore if OOS L30d>0 (F38c) or sellable WOS<4 (F38d). Buyability='Not Buyable' or "
         "'Suppressed' → W1-W4=0 + W5 catch-up = baseline × 2.50 (F38f, 25%/wk decay).",
         "Amazon forecasts now track real consumer demand. Offline ASINs get a recovery curve: zero "
         "while down, then a single catch-up order when the listing comes back, then normal cadence."),
     ]),
    ("Math fine-tuning — phantom dedupe & spike attenuation (F39-F48)",
     [
        ("F39", "Strip duplicate-order runs from history",
         "When the same large order qty appears in 3+ adjacent weeks within ±5%, it's almost always "
         "one real order being re-broadcast by upstream feeds — a phantom replay, not real demand.",
         "Detect runs of near-identical large orders and strip the duplicates from history before "
         "model classification, so the L13/L26 baselines reflect actual demand, not feed echoes.",
         "Phantom orders no longer inflate the recent quarter and trick the model into thinking "
         "demand is heavier than it really is."),
        ("F40", "Order-rate deceleration scaling",
         "After F39 dedup, sometimes the surviving orders show a clear ramp-down (e.g. last 3 nz "
         "orders 60 / 120 / 360). The forward forecast inherited the older heavier rate.",
         "If the last 3 non-zero orders are decelerating, scale the forward forecast to inherit "
         "the slowing pace instead of the heavier earlier pace. Skip when F36/F38 are already "
         "governing W1-W4 zeroing — those rules win.",
         "Forecast cadence matches the live order-rate trajectory, not a rate the customer has "
         "clearly already moved off of."),
        ("F41", "Shipment-confirmed phantom dedupe (Amazon ship-lag)",
         "Amazon's ship-lag can make one shipment surface as a near-duplicate order one week "
         "earlier or later (e.g. ship 14,328 in LW_16 then a 14,184 'order' in LW_15, ~1% drift).",
         "Use shipment evidence — the strongest signal — to dedupe at ±15% tolerance (vs F39's "
         "±5%). Runs BEFORE F39 so phantoms are stripped before any baseline math runs.",
         "FF7618 case: 86,580 / 26w → 17,760 / 26w. Phantom-replay shipments no longer compound "
         "with their underlying order to double-count demand."),
        ("F42", "POS-anchored Heuristic-baseline cap (Amazon)",
         "Even after F41/F39 strip phantoms, lumpy items on Amazon Heuristic could still anchor "
         "on a window inflated by a few surviving big orders, projecting far above POS reality.",
         "Amazon-only. Cap the Heuristic 26w avg to POS sell-through × 26 when POS is reliable. "
         "Skip when F36/F38 are already zeroing W1-W4 (they win — F42 would over-correct).",
         "SF8169 case: ~87,912 → POS-realistic ~6,500 (250/wk × 26). Heuristic projections stay "
         "anchored to what Amazon shoppers are actually buying."),
        ("F43", "Cap recent spikes to 2.0× median",
         "A single 5×-median outlier in the last 4 weeks isn't a real demand signal — it's an "
         "artifact (one-time buy, replenishment shift, etc.) that drags forward forecasts up.",
         "Window-limited: only attenuate spikes inside the recent 4w window. Cap each in-place "
         "to 2.0× the L13 median. Annotate alert text with which historical index was capped "
         "and the original value, so reviewers can audit the cap.",
         "FF25895 case: 48,816 → 64,080 / 26w with smoother profile. Recent freak orders no "
         "longer leak forward 6 months."),
        ("F44", "Re-classify after F43 attenuation",
         "Once F43 caps the recent spike, the model would re-read the now-flat recent window "
         "as 'lumpy / zero-heavy' and route to the wrong recipe — Lumpy when Dense is correct.",
         "Re-run classification on the F43-cleaned history. The dense order pattern that "
         "survives the cap usually routes to a baseline-style recipe, matching planner intent.",
         "The cap doesn't accidentally flip the model into the wrong recipe and undo its own "
         "smoothing."),
        ("F45", "Per-week forecast cap — defensive guardrail",
         "Even after F43+F44, seasonal positioning can amplify a surviving spike at one position "
         "past 2× the post-cap baseline.",
         "Cap each forecast week individually so no single position can run away from the "
         "post-cap baseline. Records the count of clamped weeks in alert text.",
         "Final-stage safety net: no single week can blow past the smoothed baseline, even when "
         "earlier rules pass it through."),
        ("F46", "Post-F44 forecast rebuild",
         "When F43+F44+F45 fire alone, the result can still be a choppy 683/wk avg with "
         "mostly-zero weeks left over from the spike attenuation — not the smooth cadence the "
         "planner expects.",
         "When F44 fires, F46 rebuilds the 26-week forecast from the customer's PRE-disruption "
         "baseline: L26 nz-mean (post-F43, post-F39, post-F41) × seasonal shape, distributed "
         "smoothly across all weeks.",
         "FF25895 case: smooth ~1,380/wk × 26 ≈ 35,880 — matches the manual planner's intended "
         "steady cadence instead of bursty zero-heavy output."),
        ("F47", "OOS rebuild-ramp normalization (2026-05-07)",
         "After a multi-week stockout, customers often place compounding catch-up orders to "
         "rebuild safety stock. These show up as 3-5× the pre-OOS pace and inflate the L13 "
         "baseline for ~6 months. F47 detects this pattern and normalizes the surge to true "
         "demand intent before downstream rules see it.",
         "Detect ≥3 consecutive ship=0 weeks with ord>0 (active stockout). Cap each within-gap "
         "order at 1.3× the pre-OOS baseline (L13 nz-avg from pre-OOS window). First post-gap "
         "week also capped if it lands ≥1.5× baseline. Capped indices passed to F39/F41 "
         "dedupe rules as `protected_indices` to prevent double-zeroing.",
         "FF12660 case: pre-OOS ~1,800/wk, then 8,640 / 9,720 / 6,840 rebuild orders. F47 "
         "caps each to ~2,340 before downstream math runs. Brought FF12660 from +27% over "
         "manual to -13% under, matching VP's pre-OOS pace expectation."),
        ("F48", "Post-OOS spike-and-cooldown anchor (2026-05-07)",
         "Even after F47 normalizes the stockout pattern, some records still get inflated L13 "
         "from a single big rebuild order followed by L4 cooling toward true pace. F47 "
         "requires ≥3 CONSECUTIVE ship=0 weeks; F48 catches the spike-and-cooldown pattern "
         "where the gap was scattered (lumpy zeros, not strictly consecutive).",
         "Trigger A (universal): max ord in L13 ≥ 2.5× median (excl max), spike in W-12..W-5, "
         "AND L4 all-weeks avg < L13 nz-avg × 0.80. Trigger B (Amazon): healthy POS AND L4 ord "
         "< POS_blend × 0.85. Action: cap baseline at MAX(L4_avg, L26_avg) × 1.20 (Amazon: "
         "MAX(L4, POS_blend) × 1.20). Critical: uses L4 ALL-weeks avg (zeros = real buyer-side "
         "pause signal), not L4 nz-avg.",
         "BB13437 (Amazon) +54% over manual → +3% (anchored on POS_blend). FF15592 (Walmart) "
         "+33% → +24% (anchored on max L4/L26). 147 fires across 4,194-record cohort."),
     ]),
    ("Bug fixes — root-cause corrections (2026-05-07)",
     [
        ("VP-Q6", "Fill-rate unit-mismatch fix",
         "The OOS detection compared `Shpd_Wk_L13W_cust_` (a per-week average) against L13 "
         "ORDER total — off by a factor of 13. False-positived OOS for nearly every active "
         "record (e.g. BB13437 reported 9% fill-rate when true fill-rate was 120% catch-up). "
         "Cascade then forced L13 nz-avg baseline + enabled F13 drawdown lift + F38 trend "
         "lift, inflating cap_base by 50-100% across thousands of records.",
         "Compare per-week-avg shipped vs per-week-avg ordered (apples-to-apples). Raised "
         "OOS threshold from 0.85 → 0.70 to require a clearly broken fill-rate. The fix "
         "removed the false-positive cascade for the majority of the active cohort.",
         "BB13437 baseline drop ~85% → matches POS reality. Cohort total dropped from prior "
         "+18% bias to -24% under manual — closer to the truth, since manual planners are "
         "themselves running ~20% over real demand based on over-projection patterns."),
        ("F36 (rev)", "F36 widening + active-orders guard",
         "F36 (Amazon stockup-burnoff) only looked at the last 13 weeks of shipments — missed "
         "stockups older than 13w (e.g. SF8169's 28k-unit stockup at W-11 was outside the "
         "narrower window). Also fired even when L4 orders were ramping (FF12853 had L4 "
         "ord 2,600/wk vs POS 1,700/wk — actively replenishing — but F36 still zeroed W1-W5).",
         "(1) Widen SHP window from L13W → L26W. (2) Add L4 active-orders guard — if recent "
         "ORDERS run at ≥70% of POS rate, customer is actively replenishing; F36 doesn't fire "
         "even if a recent shipment was big.",
         "SF8169: now correctly forced to 0 for 26w (was 8,280). FF12853: W1-W5 no longer "
         "zeroed; AI 41,316 vs prior 34,380, closer to the 51,755 manual."),
        ("R9 (rev)", "R9 unconditional Heuristic ceiling",
         "R9's L52-avg × 2.0 ceiling on Heuristic baseline was previously SKIPPED when F23b "
         "trailing-zero discount also fired, on the theory that double-discount would crush "
         "items recovering from a dip. But this let single-PO patterns escape the ceiling: "
         "FF7612 had L13 = single 5,208-unit PO 12w ago then dormant; F23b discounted to "
         "1,562 which projected ~37K vs L26 all-weeks avg 426/wk — still ~9× the truth.",
         "Apply R9 unconditionally. Multiplier 2.0× normally, raised to 2.5× when F23b "
         "also fires (softens the double-discount on items legitimately recovering).",
         "FF7612 case: 37,584 (+657% over manual) → 8,736 (+76%). Single-PO patterns no "
         "longer escape the absolute L52 ceiling."),
        ("R1 Amazon gate", "Amazon items NEVER classified as OTB",
         "R1 (One-Time-Buy detector) was designed for off-price retailers (Burlington, Ross) "
         "that buy in occasional lumps with no replenishment. Amazon items with sparse-looking "
         "histories were falsely classified as OTB — but Amazon ordering is centrally managed "
         "and even sparse histories reflect ongoing replenishment, not one-time buys.",
         "Add unconditional Amazon gate at the top of R1: if customer name contains 'AMAZON', "
         "return False from `_detect_otb` regardless of pattern. Amazon items that look "
         "OTB-shaped route to the standard Inactive recipe instead, which restarts forecasting "
         "the moment orders resume.",
         "47 Amazon records previously falsely classified as OTB (zero forecast) now route "
         "through Inactive with a real forecast path."),
     ]),
    ("Math fine-tuning — status & classification",
     [
        ("F5", "Respect 'Launching/New/Pilot' status",
         "Items flagged 'Launching/New/Pilot' in PT_Item_Status were getting routed to Inactive when their "
         "recent quarter showed zero (because they hadn't started shipping yet).",
         "If status contains LAUNCH/NEW/PILOT, skip Inactive — route through the new-item logic that uses "
         "sibling SKUs + customer baseline + ramp curves.",
         "Launching items inherit reasonable forecasts from their cousins instead of zero."),
        ("F8", "Check more category fields",
         "Seasonality was only matched on Master_Category — Product_Category and Product_Subcategory were ignored.",
         "Match on all three category fields. More items get the right seasonal shape applied.",
         "Better seasonal coverage across the catalog."),
        ("F19", "Optional: floor for Inactive items with big planner numbers",
         "Items classified Inactive but with manual projections of 5,000+ AND consumer POS > 0 were forced to zero.",
         "Optional --conservative-inactive flag: apply a 50% manual-shaped floor instead of zero for these items.",
         "Opt-in safety net for accounts where planner manuals are very high-confidence."),
        ("F20", "Respect explicit planner zeros",
         "The Reactivating recipe could produce forecasts from history even when the planner had explicitly "
         "zeroed out all 26 weeks.",
         "If the planner's manual is zero across all 26 weeks, force the Reactivating output to zero too.",
         "When the planner says 'this is dead', the AI agrees."),
     ]),
]
def _f_group_pages(items):
    """Chunk an F-group into pages.  refinement_card v4 lays out 4 sections
    (Problem/Fix/Why/Example) inside a 2.10" card.  Cap at 3 cards per page
    at 2.10" each (3 x 2.10 + 2 x 0.05 gutter + 0.85 top = 7.25", fits 7.30")."""
    n = len(items)
    if n <= 2:
        return [(items, [Inches(2.85)] * n)]
    if n == 3:
        return [(items, [Inches(2.10)] * 3)]
    # n >= 4: chunk into pages of 3 cards each at 2.10" — last page may have 1-3
    pages = []
    for i in range(0, n, 3):
        chunk = items[i:i+3]
        h = Inches(2.10) if len(chunk) >= 2 else Inches(2.85)
        pages.append((chunk, [h] * len(chunk)))
    return pages

# Pre-compute page count per F-group for accurate page numbering
_f_group_pageinfo = [(_f_group_pages(g[1]), g[0]) for g in F_GROUPS]
_F_GROUPS_TOTAL_PAGES = sum(len(p[0]) for p in _f_group_pageinfo)

def slide_f_group_page(group_idx, total_pages_in_group, page_in_group,
                       title, page_items, page_heights, abs_page_n):
    s = prs.slides.add_slide(BLANK)
    if total_pages_in_group > 1:
        sub = (f"Math fine-tuning - Group {group_idx} of {len(F_GROUPS)}"
               f" (page {page_in_group} of {total_pages_in_group})")
    else:
        sub = f"Math fine-tuning - Group {group_idx} of {len(F_GROUPS)}"
    slide_header(s, title, sub)
    top = Inches(0.85)
    for (code, title2, before, after, why), h in zip(page_items, page_heights):
        refinement_card(s, MARGIN, top, CONTENT_W, h, code, title2,
                        before, after, why, example=EXAMPLES.get(code, ""))
        # 0.05" gutter (example is now inside the card, so gutter can shrink)
        top += h + Inches(0.05)
    page_footer(s, abs_page_n, 35)
    return s

_abs_page = 20  # F-group slides start at page 20
for gi, (pages, gtitle) in enumerate(_f_group_pageinfo):
    total = len(pages)
    for pi, (page_items, page_heights) in enumerate(pages):
        page_in_group = pi + 1
        abs_page = _abs_page
        _abs_page += 1
        slides.append(lambda gtitle=gtitle, page_items=page_items,
                              page_heights=page_heights, gi=gi, total=total,
                              page_in_group=page_in_group, abs_page=abs_page:
                      slide_f_group_page(gi+1, total, page_in_group, gtitle,
                                          page_items, page_heights, abs_page))

# ── 23. Master pack & rounding ────────────────────────────────────────────────
def slide_pack():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Case-pack rounding",
                 "AI numbers always match what we can actually ship")
    add_text(s, MARGIN, Inches(1.05), CONTENT_W, Inches(0.45),
             "The rule", size=14, bold=True, color=NAVY)
    add_bullets(s, MARGIN, Inches(1.45), CONTENT_W, Inches(2.0), [
        "Every non-zero AI forecast week is rounded to the nearest case-pack multiple",
        "We pull the case-pack from Quickbase Styles (defaults to 1 if missing)",
        "If a recipe predicts 47 units and the case-pack is 12, the AI shows 48 (4 cases)",
        "Rounding happens BEFORE the PO zero-out — confirmed POs stay confirmed POs",
    ], size=12)
    add_text(s, MARGIN, Inches(3.50), CONTENT_W, Inches(0.45),
             "Why this matters", size=14, bold=True, color=NAVY)
    add_bullets(s, MARGIN, Inches(3.90), CONTENT_W, Inches(2.0), [
        "Inventory ships in case-pack increments — fractional units are physically impossible",
        "AI numbers align with how POs actually get cut (always whole cases)",
        "Replenishment downstream doesn't have to do its own rounding",
        "You see the realistic, shippable forecast — not a math abstraction",
    ], size=12)
    page_footer(s, 28, 35)
    return s
slides.append(slide_pack)

# ── 24. AI_ALERT ──────────────────────────────────────────────────────────────
def slide_alerts():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "The AI_ALERT column — your shortcut to understanding the forecast",
                 "Every refinement that fired is named in the alert text")
    add_text(s, MARGIN, Inches(1.05), CONTENT_W, Inches(0.45),
             "When does an alert fire?", size=14, bold=True, color=NAVY)
    add_bullets(s, MARGIN, Inches(1.45), CONTENT_W, Inches(0.9), [
        "Whenever AI total differs from your manual projection by more than 5% (over the 26-week window)",
    ], size=12)
    add_text(s, MARGIN, Inches(2.50), CONTENT_W, Inches(0.45),
             "What's in the alert text", size=14, bold=True, color=NAVY)
    add_bullets(s, MARGIN, Inches(2.90), CONTENT_W, Inches(2.6), [
        "The recipe used (Steady-Demand, Lumpy, Seasonal Baseline, etc.)",
        "How 'normal demand' was defined for this item (which baseline mode)",
        "Which refinement codes fired (e.g. 'F4 thin-history widened', 'F15 order-coverage 1.32x', 'VP-Q4 zeroed 3 weeks')",
        "Seasonal/event notes (Prime Day weeks lifted, Fall Deal weeks lifted, etc.)",
        "Stockout flags if VP-Q2 fired ('clean orders 18% higher than raw — chronic stockouts adjusted')",
    ], size=12)
    add_text(s, MARGIN, Inches(5.70), CONTENT_W, Inches(0.45),
             "How to use it", size=14, bold=True, color=NAVY)
    add_text(s, MARGIN, Inches(6.10), CONTENT_W, Inches(0.85),
             "Read the alert FIRST. It tells you exactly which rules drove the AI number, in plain "
             "language. If you disagree with what fired, the viewer has 'Use AI', 'Use Sugg', and "
             "manual-edit buttons to override on the spot.",
             size=12, color=GREY_DK)
    page_footer(s, 29, 35)
    return s
slides.append(slide_alerts)

# ── 25. Validation mode ───────────────────────────────────────────────────────
def slide_validate():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Validation mode — check YOUR projections",
                 "Run with --validate to flag suspicious manual numbers (no Quickbase writes)")
    rows = [
        ("CRITICAL", "Spike",         "A single week is more than 5x what we'd expect"),
        ("CRITICAL", "Demand on dead item","Non-zero projection on an item with no orders in 6+ months"),
        ("WARNING",  "Out-of-band",   "Outside the expected range (less than 30% or more than 200% of pace)"),
        ("WARNING",  "Sudden stop",   "Manual drops to zero with no decline signal in the data"),
        ("WARNING",  "Off-cadence",   "Non-zero in a week the customer normally skips (monthly orderers)"),
        ("WARNING",  "Pack mismatch", "Quantity isn't a clean case-pack multiple"),
        ("LOW",      "Variance",      "Within range but more than 5% off the AI's expected pace"),
    ]
    add_table(s, MARGIN, Inches(1.05), CONTENT_W, Inches(4.5),
              ["Severity", "Flag type", "What triggers it"], rows, size=11, hdr_size=12)
    add_text(s, MARGIN, Inches(5.80), CONTENT_W, Inches(0.45),
             "Notes", size=14, bold=True, color=NAVY)
    add_bullets(s, MARGIN, Inches(6.20), CONTENT_W, Inches(0.9), [
        "Validation reuses the same data and recipe-picker as forecast mode — so if the AI thinks something's "
        "off, validation will too",
        "Output: validation_results.json, opens in the viewer in 'validate mode'",
        "You can change the band tightness with --threshold (default 2.0; use 3.0 for more relaxed)",
    ], size=11)
    page_footer(s, 30, 35)
    return s
slides.append(slide_validate)

# ── 26. Configuration knobs ───────────────────────────────────────────────────
def slide_config():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "What admins can adjust", "Knobs and switches in plain English")
    rows = [
        ("Prime Day lift schedule", "Weeks 5–9 ramped (peak 25% in W7-W8)",
         "How much to lift Amazon orders during Prime Day pre-buy"),
        ("Fall Deal weeks / lift",  "Weeks 23–25 at +12%",
         "When and how much to lift for the September pre-order"),
        ("Amazon detector",         "Customer name contains 'AMAZON'",
         "Decides which records get Amazon-only treatment (POS pulls + Prime Day lifts)"),
        ("Steady-recipe smoothing", "0.3 / 0.1",
         "How fast the recipe reacts to recent changes vs. holds the long-run trend"),
        ("Lumpy-recipe smoothing",  "0.3",
         "How fast the lumpy recipe reacts to recent order sizes"),
        ("Seasonal dampening",      "10% (keeps shape within ±20% of normal)",
         "How much we soften time-of-year swings to prevent runaway forecasts"),
        ("Alert threshold",         "5% variance",
         "How far AI vs. manual must differ before AI_ALERT fires"),
        ("Conservative inactive",   "Off (opt-in)",
         "Whether Inactive items with big planner numbers get a small floor instead of zero"),
    ]
    add_table(s, MARGIN, Inches(1.05), CONTENT_W, Inches(4.4),
              ["Knob", "Default", "What it controls"], rows, size=11, hdr_size=12)
    add_text(s, MARGIN, Inches(5.70), CONTENT_W, Inches(0.45),
             "Command-line switches you might use", size=14, bold=True, color=NAVY)
    add_bullets(s, MARGIN, Inches(6.10), CONTENT_W, Inches(1.0), [
        "--oos-smoothing — turn on stockout-aware demand (currently off by default)",
        "--no-po-zero — turn OFF the confirmed-PO zero-out (only for back-testing)",
        "VP-Q1 (baseline definition) and VP-Q3 (cadence) are always on",
    ], size=11)
    page_footer(s, 31, 35)
    return s
slides.append(slide_config)

# ── 27. Quickbase schema ──────────────────────────────────────────────────────
def slide_qb():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "What the tool reads from / writes to Quickbase",
                 "So you know where the numbers come from")
    rows = [
        ("Projections table", "READ", "Pulls active records, your manual projections, AI columns, flag status"),
        ("Projections table", "WRITE", "Updates AI_PRJ_W1 through AI_PRJ_W26 + AI_ALERT for each record"),
        ("Styles table",      "READ", "Pulls case-pack size, brand, category fields"),
        ("Amazon Catalog",    "READ", "Pulls consumer-sales rates (Amazon items only)"),
        ("Order History (saved report 27)", "READ", "Pulls confirmed open POs by customer ship-by date"),
        ("Comments table",    "WRITE", "Saves planner comments + flags (Date-of-Week is QB-computed, not set by us)"),
    ]
    add_table(s, MARGIN, Inches(1.05), CONTENT_W, Inches(4.4),
              ["Table", "Direction", "What we use it for"], rows, size=10, hdr_size=11)
    add_text(s, MARGIN, Inches(5.70), CONTENT_W, Inches(0.45),
             "How write-back works", size=14, bold=True, color=NAVY)
    add_bullets(s, MARGIN, Inches(6.10), CONTENT_W, Inches(1.0), [
        "Writes happen in parallel (default 6 records at a time)",
        "Progress saved every 50 records — if interrupted, --resume picks up where it left off",
        "Failed writes are logged separately so you can retry without redoing successful ones",
    ], size=11)
    page_footer(s, 32, 35)
    return s
slides.append(slide_qb)

# ── 28. Reading an alert — example ────────────────────────────────────────────
def slide_example():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Example — how to read an alert in 30 seconds",
                 "Real-style alert with a step-by-step read")
    box_top = Inches(1.10)
    add_rect(s, MARGIN, box_top, CONTENT_W, Inches(2.0), fill=GREY_LT)
    add_text(s, MARGIN + Inches(0.20), box_top + Inches(0.10),
             CONTENT_W - Inches(0.4), Inches(0.40),
             "Record: 1864-FF8654   ·   Customer: Amazon   ·   Brand: Glad for Pets",
             size=11, bold=True, color=NAVY)
    add_text(s, MARGIN + Inches(0.20), box_top + Inches(0.50),
             CONTENT_W - Inches(0.4), Inches(0.40),
             "Your manual: 14,200    ·    AI says: 12,850    ·    Variance: −9.5% (alert)",
             size=11, color=GREY_DK)
    add_text(s, MARGIN + Inches(0.20), box_top + Inches(0.90),
             CONTENT_W - Inches(0.4), Inches(1.0),
             "AI_ALERT text:\n"
             "Recipe: Seasonal Baseline · Baseline: average of last 13 active weeks (486/wk) · "
             "Prime Day W5-W9 ramp applied · Order rate runs 32% over Amazon shopper rate, baseline "
             "blended toward orders · 4 trailing zero weeks → −12% drawdown discount · 2 confirmed POs "
             "in W3 and W5 totaling 1,940 units, AI weeks zeroed.",
             size=11, color=GREY_DK)
    add_text(s, MARGIN, Inches(3.30), CONTENT_W, Inches(0.45),
             "How to read this in 30 seconds", size=14, bold=True, color=NAVY)
    add_bullets(s, MARGIN, Inches(3.70), CONTENT_W, Inches(3.5), [
        "Recipe is Seasonal Baseline — this is a steady, dense item (the most common case)",
        "Baseline of 486/week comes from 'active weeks only' — not skewed by recent quiet weeks",
        "Prime Day lift was applied to W5-W9 because customer name contains AMAZON",
        "Your orders are 32% over consumer sales rate — buyer's choice for safety stock, AI respects that",
        "Last 4 weeks were zero → drawdown is happening → baseline trimmed 12%",
        "Customer already has confirmed POs in 2 weeks → AI didn't double-count",
        "Net result: AI is 9.5% under your manual. The alert tells you exactly why so you can decide.",
    ], size=12)
    page_footer(s, 33, 35)
    return s
slides.append(slide_example)

# ── 29. Refinement summary table ─────────────────────────────────────────────
def slide_summary():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Quick-lookup: refinements on one page",
                 "Match a code in your alert to its plain-English meaning")
    col_w = Inches(4.10)
    headers_t = ["Big-picture rules (4)", "Customer-type overrides (8)", "Math fine-tuning (17)"]
    cols = [
        [
            "VP-Q1  Active-weeks baseline",
            "VP-Q2  Treat stockouts as demand",
            "VP-Q3  Stop forcing fake bi-weekly",
            "VP-Q4  Don't double-count POs",
        ],
        [
            "R1  Off-price = lump buyers",
            "R2  Sparse-account 26w cap",
            "R3  Inactive item floor",
            "R5  International = 26w threshold",
            "R6  Lumpy recipe lift for big items",
            "R7  Amazon Fall Deal insert",
            "R8  Median anchor for burst items",
            "R9  Reactivating-recipe ceiling",
        ],
        [
            "F4   Thin recent history → year fallback",
            "F5   'Launching/New/Pilot' status respect",
            "F6   Recent slowdown / pickup detection",
            "F7   Anchor on seasonal peak",
            "F8   More category match fields",
            "F9   Big-volume sparse → MAX baseline",
            "F10  Decline detection (YoY-gated)",
            "F11  Prime Day ramp/taper shape",
            "F13  Drawdown-and-refill (Amazon)",
            "F14  Trust consumer sales over buyer pause",
            "F15  Order-coverage premium",
            "F16  Don't soften strong seasonal",
            "F17  Sparse cadence W1 seed",
            "F18  Lumpy z lift from Amazon POS",
            "F22  Trailing-zero discount",
            "F23  Reactivating drawdown + dampening",
            "F24  Final-baseline ceiling",
            "F35  Stockout backlog strip",
            "F36  Amazon stock-up burn-off",
            "F37  Forward inventory shortfall",
            "F38  Amazon POS-trend sensitivity",
            "F39-F41  Phantom-order dedupe",
            "F42  POS-anchored Heuristic cap",
            "F43-F46  Spike attenuation + rebuild",
            "+ F25 / F28-32 / T4 supporting rules",
        ],
    ]
    for i, (htxt, items) in enumerate(zip(headers_t, cols)):
        x = MARGIN + Inches(i * 4.20)
        add_rect(s, x, Inches(1.10), col_w, Inches(0.45),
                 fill=[NAVY, ORANGE, TEAL][i])
        add_text(s, x, Inches(1.10), col_w, Inches(0.45),
                 htxt, size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_bullets(s, x + Inches(0.10), Inches(1.65),
                    col_w - Inches(0.20), Inches(5.2),
                    items, size=10,
                    bullet_color=[NAVY, ORANGE, TEAL][i])
    page_footer(s, 34, 39)
    return s
slides.append(slide_summary)

# ── 35. How forecasts get to you (automated nightly run) ─────────────────────
def slide_workflow():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "How forecasts get to you — automated nightly",
                 "You don't run anything. The pipeline runs itself.")
    steps = [
        ("Overnight", NAVY,
         "Claude Nightly runs the forecaster automatically",
         "Pulls fresh Quickbase data → runs all 4 forecast phases → writes AI_PRJ_W1–W26 + AI_ALERT "
         "back to Quickbase. No manual command needed."),
        ("Morning", TEAL,
         "You open Quickbase like normal",
         "The Projections table now has fresh AI numbers and alert text on every active record. "
         "Records that need attention have severity tags and flagged-row tints already applied."),
        ("Review", ORANGE,
         "Open the Inventory Forecaster code page in QB",
         "Same browser, same QB login. Filter by your accounts/brands. Click any row to expand the "
         "26-week side-by-side. Read the alert text — it names exactly which rules fired."),
        ("Action", GOLD,
         "Override what doesn't look right",
         "'Use AI', 'Use Suggested', or edit any week directly in the table → 'Save All' commits "
         "your edits back to Quickbase. Flag/comment for the next planning cycle."),
    ]
    box_h = Inches(1.30); box_w = Inches(11.50); top_y = Inches(1.10)
    for i, (head, c, lead, body) in enumerate(steps):
        top = top_y + Inches(i * 1.45)
        add_rect(s, MARGIN + Inches(0.5), top, box_w, box_h, fill=WHITE, line=GREY_LT)
        add_rect(s, MARGIN + Inches(0.5), top, Inches(2.0), box_h, fill=c)
        add_text(s, MARGIN + Inches(0.5), top, Inches(2.0), box_h,
                 head, size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, MARGIN + Inches(2.6), top + Inches(0.10),
                 box_w - Inches(2.2), Inches(0.32),
                 lead, size=12, bold=True, color=NAVY)
        add_text(s, MARGIN + Inches(2.6), top + Inches(0.45),
                 box_w - Inches(2.2), box_h - Inches(0.55),
                 body, size=11, color=GREY_DK)
    page_footer(s, 35, 39)
    return s
slides.append(slide_workflow)

# ── 36. Viewer mock-up — what you'll see in the QB code page ──────────────────
def slide_viewer_mockup():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "What the viewer looks like",
                 "Inventory Forecaster code page — opens inside Quickbase, uses your QB login")
    # ── Mock browser frame ──
    frame_x, frame_y = MARGIN, Inches(1.05)
    frame_w, frame_h = CONTENT_W, Inches(5.6)
    add_rect(s, frame_x, frame_y, frame_w, frame_h, fill=WHITE, line=GREY_MD)
    # Browser title bar
    add_rect(s, frame_x, frame_y, frame_w, Inches(0.32), fill=GREY_LT)
    # Traffic-light dots
    for i, c in enumerate([RGBColor(0xFF,0x5F,0x57), RGBColor(0xFE,0xBC,0x2E), RGBColor(0x28,0xC8,0x40)]):
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  frame_x + Inches(0.10 + i*0.20),
                                  frame_y + Inches(0.08), Inches(0.14), Inches(0.14))
        dot.fill.solid(); dot.fill.fore_color.rgb = c; dot.line.fill.background()
    add_text(s, frame_x + Inches(0.85), frame_y, Inches(8), Inches(0.32),
             "pim.quickbase.com  ·  Inventory Forecaster", size=10, color=GREY_DK,
             anchor=MSO_ANCHOR.MIDDLE)

    # Topbar with filter chips + Save All
    topbar_y = frame_y + Inches(0.40)
    add_rect(s, frame_x + Inches(0.10), topbar_y, frame_w - Inches(0.20), Inches(0.40),
             fill=NAVY)
    add_text(s, frame_x + Inches(0.20), topbar_y, Inches(5), Inches(0.40),
             "Severity: All  ·  Brand: All  ·  Customer: AMAZON.COM.KYDC  ·  4,205 records",
             size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # Save buttons
    btn_x = frame_x + frame_w - Inches(2.30)
    add_rect(s, btn_x, topbar_y + Inches(0.06), Inches(0.95), Inches(0.28), fill=GREEN)
    add_text(s, btn_x, topbar_y + Inches(0.06), Inches(0.95), Inches(0.28),
             "Save All", size=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, btn_x + Inches(1.05), topbar_y + Inches(0.06),
             Inches(1.10), Inches(0.28), fill=GREY_LT)
    add_text(s, btn_x + Inches(1.05), topbar_y + Inches(0.06), Inches(1.10), Inches(0.28),
             "Export CSV", size=10, color=GREY_DK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Table header
    tbl_x = frame_x + Inches(0.10)
    tbl_y = topbar_y + Inches(0.50)
    tbl_w = frame_w - Inches(0.20)
    headers = ["Key", "Brand", "Customer", "MStyle", "Pri", "Manual", "AI", "Sugg", "AI vs Manual", "L13"]
    col_widths = [1.45, 1.10, 1.40, 0.80, 0.55, 0.75, 0.75, 0.75, 1.20, 0.75]
    # Normalize widths to total tbl_w
    total_w = sum(col_widths)
    col_widths_in = [w * (tbl_w/Inches(1)) / total_w for w in col_widths]

    add_rect(s, tbl_x, tbl_y, tbl_w, Inches(0.30), fill=NAVY)
    cur_x = tbl_x
    for h, w in zip(headers, col_widths_in):
        add_text(s, cur_x, tbl_y, Inches(w), Inches(0.30),
                 h, size=9, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cur_x += Inches(w)

    # Mock rows
    rows = [
        # key, brand, cust, mstyle, pri, manual, ai, sugg, var, l13, fill
        ("1864-FF8654",  "Glad for Pets",  "AMAZON",   "FF8654",  "MED", "14,200", "12,850", "13,400", "−9.5%",  "486", WHITE,        "warn"),
        ("1864-BB13437", "Burts Bees",     "AMAZON",   "BB13437", "LOW", "8,400",  "8,180",  "8,300",  "−2.6%",  "315", GREY_LT,      "ok"),
        ("1864-FF4934",  "Glad for Pets",  "AMAZON",   "FF4934",  "CRIT","2,100",  "5,050",  "3,900",  "+140.5%","194", RGBColor(0xFF,0xEB,0xEE), "crit"),
        ("1745-FF4771",  "Glad for Pets",  "WALMART",  "FF4771",  "MED", "13,500", "13,200", "13,400", "−2.2%",  "520", WHITE,        "ok"),
        ("1864-BB31553", "Burts Bees",     "AMAZON",   "BB31553", "LOW", "1,800",  "1,820",  "1,810",  "+1.1%",  "70",  GREY_LT,      "ok"),
        ("3022-FF8500",  "Vibrant Life",   "PETBARN",  "FF8500",  "MED", "5,500",  "4,200",  "5,100",  "−23.6%", "162", WHITE,        "warn"),
        ("1864-FF8654PM","Glad for Pets",  "AMAZON",   "FF8654",  "LOW", "0",      "0",      "0",      "0%",     "0",   GREY_LT,      "inactive"),
    ]
    row_h = Inches(0.36)
    for i, row in enumerate(rows):
        y = tbl_y + Inches(0.30) + Inches(i * 0.38)
        bg = row[10]
        # Row background
        add_rect(s, tbl_x, y, tbl_w, row_h, fill=bg, line=GREY_LT)
        # severity left rail
        sev_color = {"crit": RED, "warn": GOLD, "ok": GREEN, "inactive": GREY_MD}[row[11]]
        add_rect(s, tbl_x, y, Inches(0.04), row_h, fill=sev_color)
        # Cells
        cur_x = tbl_x
        for j, (val, w) in enumerate(zip(row[:10], col_widths_in)):
            color = GREY_DK
            bold = False
            if j == 4:  # priority pill
                pri_color = {"CRIT": RED, "MED": ORANGE, "LOW": GREEN}[val]
                add_rect(s, cur_x + Inches(0.10), y + Inches(0.07),
                         Inches(w - 0.20), Inches(0.22), fill=pri_color)
                add_text(s, cur_x, y, Inches(w), row_h,
                         val, size=8, bold=True, color=WHITE,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            elif j == 8:  # variance — color by direction
                if "+" in val and val != "0%":
                    color = GREEN
                elif "−" in val and val != "−0%":
                    color = RED
                bold = True
                add_text(s, cur_x, y, Inches(w), row_h,
                         val, size=9, bold=bold, color=color,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            else:
                align = PP_ALIGN.RIGHT if j >= 5 else PP_ALIGN.LEFT
                if j == 6:
                    color = NAVY; bold = True
                add_text(s, cur_x + Inches(0.05), y, Inches(w - 0.10), row_h,
                         val, size=9, bold=bold, color=color,
                         align=align, anchor=MSO_ANCHOR.MIDDLE)
            cur_x += Inches(w)

    # ── Legend below the frame ──
    leg_y = frame_y + frame_h + Inches(0.10)
    legend_items = [
        (RED,    "CRITICAL — variance >5x or alert fired"),
        (GOLD,   "WARNING — variance 2x-5x"),
        (GREEN,  "LOW — within band"),
        (RGBColor(0xFF,0xEB,0xEE), "Pink wash — flagged for follow-up"),
    ]
    cur_x = MARGIN
    for col, txt in legend_items:
        add_rect(s, cur_x, leg_y + Inches(0.05), Inches(0.18), Inches(0.18), fill=col)
        add_text(s, cur_x + Inches(0.25), leg_y, Inches(3.0), Inches(0.30),
                 txt, size=9, color=GREY_DK, anchor=MSO_ANCHOR.MIDDLE)
        cur_x += Inches(3.20)
    page_footer(s, 36, 39)
    return s
slides.append(slide_viewer_mockup)

# ── 37. Viewer enhancements — recent additions for planners ──────────────────
def slide_viewer_enhancements():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Viewer enhancements — what's new in the QB code page",
                 "Quality-of-life upgrades aimed at planner triage speed")
    cards = [
        ("L4W column", NAVY,
         "New \"Ord/Wk L4W\" column",
         "A 4-week order-rate column now sits immediately to the left of \"Ord/Wk L13W\". "
         "Planners can scan L4W vs L13W side-by-side to spot acceleration or deceleration "
         "without expanding the row. The L4W signal is what drives F6/F26/F27 baseline lifts "
         "and cuts — surfacing it on the summary row makes those rule firings self-explanatory."),
        ("LY actuals rows", TEAL,
         "Ordered LY · Shipped LY rows in the detail pane",
         "Open any row's 26-week side-by-side and you now see two extra rows below \"Suggested\": "
         "\"Ordered LY\" (green) and \"Shipped LY\" (blue), aligned to W1-W26 by pulling the values "
         "from 52 weeks before each forecast week (Ord_LW-26..LW-51 and Shp_LW-26..LW-51). "
         "Sanity-check seasonal lifts and event-week shapes against last-year reality without "
         "leaving the row. Coverage on May 7: 2,691 records have non-zero LY ord; 2,623 have non-zero LY shp."),
        ("Sticky column headers", ORANGE,
         "Header row freezes during scroll",
         "On long lists, the column headers (Manual / AI / Suggested / variance / L13 / L4W / etc.) "
         "now stay pinned at the top of the table while you scroll. CSS fix: border-collapse:separate "
         "+ box-shadow replaces border-bottom so the sticky position holds. No more lost context "
         "after the first 20 rows."),
        ("Wider CSV exports", GOLD,
         "Both export buttons now produce 160-column files",
         "\"Export All In View\" and \"Export Flagged\" both produce 160-column wide CSVs: 28 summary "
         "columns + 5 series × 26 weeks (Manual / AI / Suggested / Ordered LY / Shipped LY). "
         "Reviewers get the full week-level picture for offline analysis, audit trails, or "
         "side-by-side compares in Excel."),
    ]
    box_h = Inches(1.30); box_w = Inches(11.50); top_y = Inches(1.10)
    for i, (head, c, lead, body) in enumerate(cards):
        top = top_y + Inches(i * 1.45)
        add_rect(s, MARGIN + Inches(0.5), top, box_w, box_h, fill=WHITE, line=GREY_LT)
        add_rect(s, MARGIN + Inches(0.5), top, Inches(2.0), box_h, fill=c)
        add_text(s, MARGIN + Inches(0.5), top, Inches(2.0), box_h,
                 head, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, MARGIN + Inches(2.6), top + Inches(0.10),
                 box_w - Inches(2.2), Inches(0.32),
                 lead, size=12, bold=True, color=NAVY)
        add_text(s, MARGIN + Inches(2.6), top + Inches(0.45),
                 box_w - Inches(2.2), box_h - Inches(0.55),
                 body, size=10, color=GREY_DK)
    page_footer(s, 37, 39)
    return s
slides.append(slide_viewer_enhancements)

# ── 38. Latest run results — May 7, 2026 dataset snapshot ────────────────────
def slide_latest_run():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Latest run results — May 7, 2026",
                 "Full --all forecast against live Quickbase data")
    # Top KPI strip
    kpis = [
        ("4,194", "Active records processed", NAVY),
        ("3,526", "ALERT records (>5% var. vs manual)", ORANGE),
        ("AI_PRJ_W1..W26", "Written back to Quickbase + AI_ALERT", TEAL),
    ]
    kpi_y = Inches(1.10); kpi_h = Inches(1.10); kpi_w = Inches(4.10)
    for i, (val, lbl, c) in enumerate(kpis):
        x = MARGIN + Inches(i * 4.20)
        add_rect(s, x, kpi_y, kpi_w, kpi_h, fill=c)
        # Smaller font for the longer text in the third KPI
        val_size = 22 if len(val) <= 12 else 16
        add_text(s, x, kpi_y + Inches(0.05), kpi_w, Inches(0.55),
                 val, size=val_size, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, kpi_y + Inches(0.60), kpi_w, Inches(0.45),
                 lbl, size=11, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Model-split table
    add_text(s, MARGIN, Inches(2.40), CONTENT_W, Inches(0.35),
             "Model split across the 4,194 records", size=14, bold=True, color=NAVY)
    headers = ["Model / classification", "Records"]
    rows = [
        ("Sparse Intermittent",                 "1,166"),
        ("Croston's",                           "1,017"),
        ("Seasonal Baseline",                   "750"),
        ("Inactive",                            "402"),
        ("OTB-zero",                            "321"),
        ("Inactive (zero history)",             "278"),
        ("Heuristic",                           "143"),
        ("Reactivating",                        "46"),
        ("Pre-launch",                          "44"),
        ("Other / smaller buckets",             "27"),
    ]
    add_table(s, MARGIN, Inches(2.80), Inches(6.20), Inches(3.80),
              headers, rows, size=11, hdr_size=12)
    # LY coverage callout
    cb_x = MARGIN + Inches(6.50); cb_y = Inches(2.80); cb_w = Inches(5.80); cb_h = Inches(3.80)
    add_rect(s, cb_x, cb_y, cb_w, cb_h, fill=GREY_LT, line=GREY_LT)
    add_text(s, cb_x + Inches(0.20), cb_y + Inches(0.15), cb_w - Inches(0.40), Inches(0.35),
             "Last-year actuals coverage", size=14, bold=True, color=NAVY)
    add_bullets(s, cb_x + Inches(0.20), cb_y + Inches(0.55), cb_w - Inches(0.40), Inches(1.40), [
        "2,691 records have non-zero \"Ordered LY\" data (Ord_LW-26..LW-51)",
        "2,623 records have non-zero \"Shipped LY\" data (Shp_LW-26..LW-51)",
        "Surfaced as the new green/blue rows in the row detail pane",
    ], size=11)
    add_text(s, cb_x + Inches(0.20), cb_y + Inches(2.00), cb_w - Inches(0.40), Inches(0.30),
             "Why ALERT count is high (3,526)", size=12, bold=True, color=ORANGE)
    add_bullets(s, cb_x + Inches(0.20), cb_y + Inches(2.32), cb_w - Inches(0.40), Inches(1.40), [
        "Manual projections are still flat-rate plans for many accounts",
        "AI applies recipe + seasonal + event + recent-trend lifts per record",
        ">5% variance is intentionally a tight band — it surfaces every meaningful disagreement",
        "Use Severity + Brand filters to narrow the alert list to your own scope",
    ], size=10)
    page_footer(s, 38, 39)
    return s
slides.append(slide_latest_run)

# ── 39. Closing — your day-to-day in the QB code page ────────────────────────
def slide_close():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Your day-to-day in the QB code page",
                 "Three things you'll do; everything else is automatic")
    add_text(s, MARGIN, Inches(1.05), CONTENT_W, Inches(0.45),
             "1. Triage what needs attention", size=14, bold=True, color=NAVY)
    add_bullets(s, MARGIN + Inches(0.30), Inches(1.45), CONTENT_W - Inches(0.30),
                Inches(1.4), [
        "Filter by your accounts / brands at the top of the page",
        "Sort by severity (CRITICAL first), or just scroll — flagged rows show with a light-red wash",
        "Each row's variance % tells you at a glance how far AI is from your manual projection",
    ], size=11)
    add_text(s, MARGIN, Inches(2.95), CONTENT_W, Inches(0.45),
             "2. Read the alert before deciding", size=14, bold=True, color=NAVY)
    add_bullets(s, MARGIN + Inches(0.30), Inches(3.35), CONTENT_W - Inches(0.30),
                Inches(1.6), [
        "Click a row to expand the 26-week side-by-side: Manual · AI · Suggested · Shipments · Orders",
        "The AI_ALERT column names every rule that fired — recipe, baseline mode, refinements, PO zero-outs",
        "Most 'why is the AI saying that?' questions are answered by reading the alert text",
    ], size=11)
    add_text(s, MARGIN, Inches(5.20), CONTENT_W, Inches(0.45),
             "3. Override + flag where needed", size=14, bold=True, color=NAVY)
    add_bullets(s, MARGIN + Inches(0.30), Inches(5.60), CONTENT_W - Inches(0.30),
                Inches(1.4), [
        "'Use AI' or 'Use Suggested' to copy those numbers into your manual projection",
        "Or edit any week directly in the table → 'Save All' commits your edits back to Quickbase",
        "Add a comment + flag so the next nightly run remembers the record needs attention",
    ], size=11)
    page_footer(s, 39, 39)
    return s
slides.append(slide_close)

# ── Text-overflow QA ─────────────────────────────────────────────────────────
# Heuristic estimator: catches the common cases (textbox / table cell whose
# wrapped content cannot fit in its given height).  Conservative on purpose —
# we'd rather flag a borderline case than let an obvious overflow slip.
def _estimate_text_height_emu(text, width_emu, font_size_pt):
    """Approximate rendered height of `text` inside `width_emu` at `font_size_pt`.
    Uses an avg char width of 0.55x font size (conservative for mixed-case
    Calibri body text — wider chars like M/W push this up), 1.22x line spacing.
    Splits on '\n' first (hard breaks) then word-wraps each line."""
    if not text:
        return 0
    # Convert EMU -> points (1 in = 914400 EMU = 72 pt)
    width_pt = width_emu / 914400.0 * 72.0
    avg_char_pt = font_size_pt * 0.55  # conservative — flag earlier
    line_h_pt   = font_size_pt * 1.22  # conservative line spacing
    # Account for textbox internal margins (~0.05" L + 0.05" R = 7.2 pt total)
    usable_pt = max(width_pt - 7.2, 12.0)
    chars_per_line = max(1, int(usable_pt / avg_char_pt))
    total_lines = 0
    for hard_line in str(text).split("\n"):
        if not hard_line.strip():
            total_lines += 1
            continue
        # Word-wrap: count lines as ceil(len / chars_per_line) but smarter on words
        words = hard_line.split(" ")
        cur = 0
        line_count = 0
        for w in words:
            wl = len(w) + (1 if cur > 0 else 0)
            if cur + wl > chars_per_line and cur > 0:
                line_count += 1
                cur = len(w)
            else:
                cur += wl
        if cur > 0:
            line_count += 1
        total_lines += max(1, line_count)
    return int(total_lines * line_h_pt / 72.0 * 914400.0)

def _check_overflow(presentation):
    """Walk every textbox and table cell; warn when content likely overflows
    its allocated height.  Returns list of (slide_idx, shape_name, detail) tuples."""
    issues = []
    for s_idx, slide in enumerate(presentation.slides, 1):
        for shape in slide.shapes:
            # 1) Plain textboxes / autoshapes with text
            if shape.has_text_frame and not shape.has_table:
                tf = shape.text_frame
                # Skip empty
                if not (tf.text or "").strip():
                    continue
                # Default font size = 18 pt unless overridden per run
                # We'll use the LARGEST run size observed as the line-height proxy
                max_pt = 0
                concat_text = []
                for para in tf.paragraphs:
                    line_text = "".join(r.text for r in para.runs) or para.text
                    concat_text.append(line_text)
                    for r in para.runs:
                        if r.font.size:
                            max_pt = max(max_pt, r.font.size.pt)
                if max_pt == 0:
                    max_pt = 14  # our default
                joined = "\n".join(concat_text)
                est_h = _estimate_text_height_emu(joined, shape.width, max_pt)
                if est_h > shape.height:
                    excess = (est_h - shape.height) / 914400.0
                    if excess > 0.03:  # tiny near-fit tolerance only
                        snippet = joined.replace("\n", " / ")[:90].encode("ascii", "replace").decode("ascii")
                        issues.append((s_idx, "textbox",
                                       f"+{excess:.2f}\" over @ {max_pt:.0f}pt: \"{snippet}...\""))
            # 2) Tables — check each cell
            if shape.has_table:
                tbl = shape.table
                for r_idx, row in enumerate(tbl.rows):
                    row_h = row.height
                    for c_idx, cell in enumerate(row.cells):
                        cell_w = cell.width if hasattr(cell, "width") else shape.width // len(tbl.columns)
                        ctxt = cell.text_frame.text or ""
                        if not ctxt.strip():
                            continue
                        max_pt = 0
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.size:
                                    max_pt = max(max_pt, run.font.size.pt)
                        if max_pt == 0:
                            max_pt = 11
                        est_h = _estimate_text_height_emu(ctxt, cell_w, max_pt)
                        if est_h > row_h + 36576:
                            excess = (est_h - row_h) / 914400.0
                            if excess > 0.05:
                                snippet = ctxt.replace("\n", " / ")[:80].encode("ascii", "replace").decode("ascii")
                                issues.append((s_idx, f"table[{r_idx},{c_idx}]",
                                               f"+{excess:.2f}\" over @ {max_pt:.0f}pt: \"{snippet}…\""))
    return issues

# ── Build it ─────────────────────────────────────────────────────────────────
import sys
print(f"Building {len(slides)} slides...")
for fn in slides:
    fn()

# Pre-save overflow audit (always runs unless --allow-overflow passed)
ALLOW_OVERFLOW = "--allow-overflow" in sys.argv
issues = _check_overflow(prs)
if issues:
    print(f"\n[!] Text-overflow audit: {len(issues)} potential issue(s):")
    for s_idx, where, detail in issues[:50]:
        print(f"  - slide {s_idx} [{where}] {detail}")
    if len(issues) > 50:
        print(f"  - ...and {len(issues) - 50} more")
    if not ALLOW_OVERFLOW:
        print("\n[X] Build halted - fix overflow or rerun with --allow-overflow to override.")
        sys.exit(2)
    print("\n(continuing because --allow-overflow was passed)")
else:
    print("[OK] Text-overflow audit: no issues detected.")

prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slide count: {len(prs.slides)}")
